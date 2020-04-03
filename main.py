from PyQt5.QtWidgets import *
from PyQt5.uic import loadUi
from PyQt5.QtCore import Qt,QDateTime,QPropertyAnimation,QRect,QPoint,QEasingCurve
from PyQt5.Qt import QFont,QDate,QBrush,QColor,QObject,QPixmap
from PyQt5 import QtGui
import matplotlib.pyplot as plt
import matplotlib
from matplotlib.ticker import FuncFormatter
from matplotlib.backends.backend_qt5agg import FigureCanvas
from matplotlib.figure import Figure
import time,requests,json,re,sys,os,shutil,random,threading
from datetime import timedelta,datetime
import numpy as np
import pandas as pd
from pandas.io.json import json_normalize
try:
    from pptx import Presentation
    from pptx.util import Pt,Inches
    from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import ColorFormat, RGBColor
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.enum.action import PP_ACTION
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    import Examtest
    import win32com.client
except:pass
#import pymysql

#先檢查是否有D:/python
while os.path.isdir('D:/python')!=True:
    try:os.mkdir('D:/python')
    except:pass

def alert(word):QMessageBox.warning(w,'警告!',str(word),QMessageBox.Ok)

def search_summary():
    global rs,AuthToken,summary_all,item_all,httpm
    try:w.summary.cellChanged.disconnect()#關閉"計算使用者勾選報表table的"連接
    except:pass
    w.stackedWidget_su_it.setCurrentIndex(0)#將畫面轉到summary
    w.Search_scatter.setDisabled(1)#重新搜尋SUMMARY時,畫scatter功能就要關掉
    w.Status.setText('Status: 搜尋summary中...')
    w.summary.setSortingEnabled(False)#搜尋前>>先讓table不可排序>>因為可排序後搜尋會導致有些格子為空值
    while w.summary.rowCount() > 0: w.summary.removeRow(w.summary.rowCount()-1) #先清空TABLE
    while w.item_su.rowCount() > 0: w.item_su.removeRow(w.item_su.rowCount()-1) #先清空TABLE
    while w.item_search_su.rowCount() > 0: w.item_search_su.removeRow(w.item_search_su.rowCount()-1) #先清空TABLE
    if w.lineEdit_item.text()=='' and w.lineEdit_lot.text()=='' and w.lineEdit_file.text()=='' and w.lineEdit_pg.text()=='' and w.lineEdit_tester.text()=='':
        w.Status.setText('Status: 沒輸入任何搜尋資料，請確認...')
        alert('沒輸入任何搜尋資料，請確認...')
        return
    lotid,filename,program,tester,cmod,device,start_time,end_time=[],[],[],[],[],[''],'',''
    lotid=(str(w.lineEdit_lot.text()).replace(' ','').split(','))
    filename=(str(w.lineEdit_file.text()).replace(' ','').split(','))
    program=(str(w.lineEdit_pg.text()).replace(' ','').split(','))
    tester=(str(w.lineEdit_tester.text()).replace(' ','').split(','))
    cmod=(str(w.lineEdit_cmod.text()).replace(' ','').split(','))
    if w.Need_time.isChecked():
        start_time=w.dateEdit_1.date().toString("yyyy-MM-dd")
        end_time=w.dateEdit_2.date().toString("yyyy-MM-dd")
    w.Status.setText('Status: 登入中..')
    #------------------------登入------------------------
    #選擇哪一廠>>中華\銅鑼
    Factory = str(w.factory.currentText())
    if Factory=='中華': httpm = '******'
    else: httpm = '******'
    try_max_time=5
    try_time=0
    while try_time<try_max_time:
        try:
            rs=requests.session()
            header={
                'Key some request information
            }
            js=json.loads('need password')
            res=rs.post(httpm+'/api/login',headers=header,json=js)
            AuthToken=res.headers['AuthToken']
            break
        except:
            try_time = try_time +1
            time.sleep(try_time)
    #----------------------------------------------得到summary----------------------------------------------
    summary_count_0 = False#如果沒有搜尋到任何一筆報表則為True
    #寫入搜尋資料(json檔)
    try_time=0
    while try_time<try_max_time:
        try:
            js_search1 = '''{"PreDays":null,"EngineeringData":"YES","FileType":"FT","DefaultQuery":false,
            "poptions":{"pagesize":2000,"skipnumber":0,"gridsort":{"START_T":-1},"filter":{}}'''
            if cmod[0] !='':js_search1 = js_search1 + ',"CMOD_COD":[]'
            if filename[0] !='':js_search1 = js_search1 + ',"FILE_NAME":[]'
            if device[0] !='': js_search1 = js_search1 + ',"PART_TYP":[]'
            if lotid[0] !='': js_search1 = js_search1 + ',"LOT_ID":[]'
            if program[0] !='': js_search1 = js_search1 + ',"JOB_NAM":[]'
            if tester[0] !='': js_search1 = js_search1 + ',"NODE_NAM":[]'
            if str.strip(start_time) !='' and str.strip(start_time) !='' : js_search1 = js_search1 + ',"START_T":{"gteVal":"'+str.strip(start_time)+'","lteVal":"'+str.strip(end_time)+'"}'
            js_search1 = js_search1 + '}'
            header2={'User-Authorization':AuthToken}
            js2=json.loads(js_search1)
            if cmod[0] !='':js2['CMOD_COD']=cmod
            if filename[0] !='':js2['FILE_NAME']=filename
            if lotid[0] !='':js2['LOT_ID']=lotid
            if program[0] !='':js2['JOB_NAM']=program
            if tester[0] !='': js2['NODE_NAM']=tester
            if device[0] !='':js2['PART_TYP']=device
            res2=rs.post(httpm+'/api/datalog/getsearchresult',headers=header2,json=js2)
            fuckyouhaha = res2.json()['count']#如果這邊抓的到表示順利搜尋完成
            break
        except:
            try_time = try_time +1
            time.sleep(try_time)
    if try_time >= try_max_time:
        w.Status.setText('Status: 接收昱冠資料失敗')
        alert('接收昱冠資料失敗')
        return 
    summary_all = res2.json()
    if int(summary_all['count'])==0:
        alert('抱歉....沒搜尋到任何summary')
        w.Status.setText('Status: 抱歉....沒搜尋到任何summary')
        time.sleep(2)
    w.Status.setText('Status: 搜尋summary完成')
    for xyz in range(len(summary_all['docs'])):
        data = summary_all['docs'][xyz]
        #羅列SUMMARY
        row = w.summary.rowCount()
        w.summary.setRowCount(row + 1)
        item = QTableWidgetItem()
        item.setFlags(item.flags())
        item.setCheckState(Qt.Unchecked)
        w.summary.setItem(row,0,item)
        cool = ['LOT_ID','START_T','ROM_COD','CMOD_COD','UNITS','YIELD','NODE_NAM','FILE_NAME','ENG_ID','HAND_ID','LOAD_ID','JOB_REV','JOB_NAM']
        for coo in range(len(cool)):
            if cool[coo] =='UNITS':
                try:
                    nameItem = QTableWidgetItem()
                    nameItem.setData(Qt.DisplayRole,int(data['UNITS']))
                    w.summary.setItem(row,coo+1,nameItem)
                except:w.summary.setItem(row,coo+1,QTableWidgetItem("No Data"))
            elif cool[coo] =='YIELD':
                try:w.summary.setItem(row,coo+1,QTableWidgetItem(str(data['YIELD']*100)[:4]+'%'))
                except:w.summary.setItem(row,coo+1,QTableWidgetItem("No Data"))
            else:
                try:w.summary.setItem(row,coo+1,QTableWidgetItem(str(data[cool[coo]])))
                except:w.summary.setItem(row,coo+1,QTableWidgetItem(""))
    w.summary.setSortingEnabled(True)#搜尋後>>再讓table可排序
    w.Search_item.setEnabled(1)#搜尋完summary後就可以開啟Search_item的篩選功能了
    w.Search_scatter.setDisabled(1)#把Search_scatter關掉
    w.pushButton_Group_By.setDisabled(1)#把Group_summary關掉
    w.pushButton_CID_Search_2.setDisabled(1)#把Chip_ID關掉
    w.summary.cellChanged.connect(Summary_Count)#開啟計算使用者勾選報表table的連接
    

def summary_search():
    try:w.summary.cellChanged.disconnect()#關閉"計算使用者勾選報表table的"連接
    except:pass
    w.stackedWidget_su_it.setCurrentIndex(0)#將畫面轉到summary
    cool = ['LOT_ID','START_T','ROM_COD','CMOD_COD','UNITS','YIELD','NODE_NAM','FILE_NAME','ENG_ID','HAND_ID','LOAD_ID','JOB_REV','JOB_NAM']
    title_search=w.summary_title.currentText()#要搜尋哪個title
    before_search_filename=[]#紀錄搜尋前table的filename
    #先確定有東西可搜尋
    try:
        if len(summary_all['docs'])==0:
            w.Status.setText('Status: 沒任何summary可搜尋')
            alert('沒任何summary可搜尋')
            return
    except:
        w.Status.setText('Status: 沒任何summary可搜尋')
        alert('沒任何summary可搜尋')
        return
    w.summary.setSortingEnabled(False)#搜尋前>>先讓table不可排序>>因為可排序後搜尋會導致有些格子為空值
    #Reset list>>如同勾選掉所有的篩選
    if title_search=='清除所有篩選':
        while w.summary.rowCount() > 0: w.summary.removeRow(w.summary.rowCount()-1) #先清空TABLE
        for xyz in range(len(summary_all['docs'])):
            data = summary_all['docs'][xyz]
            row = w.summary.rowCount()
            w.summary.setRowCount(row + 1)
            item = QTableWidgetItem()
            item.setFlags(item.flags())
            item.setCheckState(Qt.Unchecked)
            w.summary.setItem(row,0,item)
            for coo in range(len(cool)):
                if cool[coo] =='UNITS':
                    try:
                        nameItem = QTableWidgetItem()
                        nameItem.setData(Qt.DisplayRole,int(data['UNITS']))
                        w.summary.setItem(row,coo+1,nameItem)
                    except:w.summary.setItem(row,coo+1,QTableWidgetItem("No Data"))
                elif cool[coo] =='YIELD':
                    try:w.summary.setItem(row,coo+1,QTableWidgetItem(str(data['YIELD']*100)[:4]+'%'))
                    except:w.summary.setItem(row,coo+1,QTableWidgetItem("No Data"))
                else:
                    try:w.summary.setItem(row,coo+1,QTableWidgetItem(str(data[cool[coo]])))
                    except:w.summary.setItem(row,coo+1,QTableWidgetItem(""))
        return
    
    #先確認目前summary table有哪些>>記錄起來>>這樣就能篩選多個條件
    for row in range(w.summary.rowCount()):before_search_filename.append(w.summary.item(row,column_summary.index('Filename')).text())
    while w.summary.rowCount() > 0: w.summary.removeRow(w.summary.rowCount()-1) #先清空TABLE
    for xyz in range(len(summary_all['docs'])):
        data = summary_all['docs'][xyz]
        #至少要是搜尋前的Filename
        if data['FILE_NAME'] not in before_search_filename:continue

        if title_search not in ['Dlog_Units','Dlog_Yield']:
            #篩選搜尋的條件
            Condition_Need = False
            Condition = (str(w.lineEdit_summary.text()).split(','))
            for pppp in Condition:
                if str(data[cool[column_summary.index(title_search)-1]]).find(pppp)!=-1:Condition_Need=True
            if Condition_Need==False:continue
        else:
            #篩選搜尋的條件['UNITS']/['YIELD']
            if title_search=='Dlog_Units':
                try:
                    fuckyouhaha = int(w.lineEdit_summary.text())#如果這邊就error表示使用者亂輸入>>不予理會
                    if w.comboBox_count.currentText()=='>=' and int(data['UNITS']) < fuckyouhaha:continue
                    if w.comboBox_count.currentText()=='<=' and int(data['UNITS']) > fuckyouhaha:continue
                except:pass
            if title_search=='Dlog_Yield' :
                try:
                    fuckyouhaha = float(w.lineEdit_summary.text())#如果這邊就error表示使用者亂輸入>>不予理會
                    if w.comboBox_count.currentText()=='>=' and float(data['YIELD']) < fuckyouhaha:continue
                    if w.comboBox_count.currentText()=='<=' and float(data['YIELD']) > fuckyouhaha:continue
                except:pass

        #羅列SUMMARY
        row = w.summary.rowCount()
        w.summary.setRowCount(row + 1)
        item = QTableWidgetItem()
        item.setFlags(item.flags())
        item.setCheckState(Qt.Unchecked)
        w.summary.setItem(row,0,item)
        for coo in range(len(cool)):
            if cool[coo] =='UNITS':
                try:w.summary.setItem(row,coo+1,QTableWidgetItem(str(data['UNITS'])))
                except:w.summary.setItem(row,coo+1,QTableWidgetItem("No Data"))
            elif cool[coo] =='YIELD':
                try:w.summary.setItem(row,coo+1,QTableWidgetItem(str(data['YIELD']*100)[:4]+'%'))
                except:w.summary.setItem(row,coo+1,QTableWidgetItem("No Data"))
            else:
                try:w.summary.setItem(row,coo+1,QTableWidgetItem(str(data[cool[coo]])))
                except:w.summary.setItem(row,coo+1,QTableWidgetItem(""))
    w.summary.setSortingEnabled(True)#搜尋後>>再讓table可排序
    w.summary.cellChanged.connect(Summary_Count)#開啟計算使用者勾選報表table的連接



def search_item():
    global rs,AuthToken,summary_need,summary_all,item_all,column,findata,findata_romcod,findata_unit,findata_tester,findata_id,findata_time,findata_TP,findata_LB,findata_diff,findata_device,findata_TP_rev,findata_Handler,httpm
    w.stackedWidget_su_it.setCurrentIndex(2)#將畫面轉到item
    w.item_su.setSortingEnabled(False)#搜尋前>>先讓table不可排序>>因為可排序後搜尋會導致有些格子為空值
    sb_need=[]#放要搜尋的sb用
    try:sb_need=[int(w.comboBox_sb.currentText())]#要的sb
    except:pass
    if summary_all==''or w.summary.rowCount()==0:
        w.Status.setText('Status: 沒點選任何summary...')
        alert('沒點選任何summary...')
        return
    w.Status.setText('Status: 搜尋item中...')
    while w.item_su.rowCount() > 0: w.item_su.removeRow(w.item_su.rowCount()-1) #先清空TABLE
    while w.item_search_su.rowCount() > 0: w.item_search_su.removeRow(w.item_search_su.rowCount()-1) #先清空TABLE
    summary_need = []
    findata=[]#FILE_NAME
    findata_romcod=[]#FILE_STEP
    findata_unit=[]#FILE_UNITS
    findata_tester=[]#FILE_NODE_NAM
    findata_id=[]#FILE_LOT_ID
    findata_time=[]#FINISH_T
    findata_TP=[]#JOB_NAM
    findata_LB=[]#LOAD_ID
    findata_diff=[]#ENG_ID
    findata_device=[]#PART_TYP
    findata_TP_rev=[]#JOB_REV
    findata_Handler=[]#HAND_ID    
    for xyz in range(len(summary_all['docs'])):
        clickedFalse = True #確認此filename是否有被勾選
        data = summary_all['docs'][xyz]
        for row in range(w.summary.rowCount()):
            try:
                if w.summary.item(row,0).checkState()==2 and w.summary.item(row,column_summary.index('Filename')).text()==data['FILE_NAME']:
                    clickedFalse = False
                    #print(data['FILE_NAME'])
                    break
            except:pass
        if clickedFalse:continue
        findata.append(data['FILE_NAME'])
        try:findata_romcod.append(data['ROM_COD'])
        except:findata_romcod.append('None')
        findata_unit.append(data['UNITS'])
        findata_tester.append(data['NODE_NAM'])
        findata_id.append(data['LOT_ID'])
        try:findata_time.append(data['FINISH_T'])
        except:findata_time.append(data['START_T'])
        findata_TP.append(data['JOB_NAM'])
        findata_LB.append(data['LOAD_ID'])
        findata_diff.append(data['ENG_ID'])
        findata_device.append(data['PART_TYP'])
        findata_TP_rev.append(data['JOB_REV'])
        findata_Handler.append(data['HAND_ID']) 
        summary_need.append(summary_all['docs'][xyz])
    if len(findata)==0:
        w.Status.setText('Status: 沒偵察有勾選summary...')
        alert('沒偵察有勾選summary...')
        return

    #===時間重新排列RIT>>Right_in_Time(def search_scatter()也有用到這個參數)
    #print(findata_time)
    RIT_df=pd.DataFrame({"a1":summary_need,   "a2":findata,
                         "a3":findata_romcod, "a4":findata_unit,
                         "a5":findata_tester, "a6":findata_id,
                         "a7":findata_time,   "a8":findata_TP,
                         "a9":findata_LB,     "a10":findata_diff,
                         "a11":findata_device,"a12":findata_TP_rev,
                         "a13":findata_Handler})
    RIT_df=RIT_df.sort_values(by=['a7'],ascending=True)
    summary_need=list(RIT_df.a1)
    findata=list(RIT_df.a2)
    findata_romcod=list(RIT_df.a3)
    findata_unit=list(RIT_df.a4)
    findata_tester=list(RIT_df.a5)
    findata_id=list(RIT_df.a6)
    findata_time=list(RIT_df.a7)
    findata_TP=list(RIT_df.a8)
    findata_LB=list(RIT_df.a9)
    findata_diff=list(RIT_df.a10)
    findata_device=list(RIT_df.a11)
    findata_TP_rev=list(RIT_df.a12)
    findata_Handler=list(RIT_df.a13)
    #print(findata_time)

    #===開始查詢資料
    #選擇哪一場>>中華\銅鑼
    Factory = str(w.factory.currentText())
    if Factory=='中華': httpm = '******'
    else: httpm = '******'
    print(httpm)
    try_max_time=5
    try_time = 0
    while try_time < try_max_time:
        try:
            rs=requests.session()
            if len(sb_need)!=0:#要指定查的sb則需要用比較久的查詢方式
                js_search2='******'
                header3={'******'}
                js3=json.loads(js_search2)
                js3['files']=findata
                js3['bins']=sb_need
                res3=rs.post(httpm+'/api/datalog/filter/test/accu',headers=header3,json=js3)
                fuckyouhaha = res3.json()[0]["FAILS"]#如果這邊就error就表示沒有抓到資料
            else:#不需要指定查的sb則用較快的查詢方式
                js_search2='******'
                header3={'******'}
                js3=json.loads(js_search2)
                js3['files']=findata
                js3['total']=sum(findata_unit)
                res3=rs.post(httpm+'/api/datalog/filter/getHighestFails',headers=header3,json=js3)
                fuckyouhaha = res3.json()[0]["FAILS"]#如果這邊就error就表示沒有抓到資料
            break
        except:
           try_time = try_time + 1
           time.sleep(try_time)
    if try_time >= try_max_time:
        w.Status.setText('Status: 接收昱冠資料失敗')
        alert('接收昱冠資料失敗')
        return        
    item_all = res3.json()
    #羅列item
    only20000 =  0 #預設為顯示前20000個item
    for it in item_all:
        item = QTableWidgetItem()
        item.setFlags(item.flags())
        item.setCheckState(Qt.Unchecked)
        row = w.item_su.rowCount()
        w.item_su.setRowCount(row + 1)
        w.item_su.setItem(row,0,item)
        for xyz in range(1,len(column)):
            try:
                if str(column[xyz])=='FAILS' or str(column[xyz])=='TEST_NUM' or str(column[xyz])=='EXECS' or str(column[xyz])=='SEQ':
                    nameItem = QTableWidgetItem()
                    nameItem.setData(Qt.DisplayRole,int(it[column[xyz]]))
                    w.item_su.setItem(row,xyz,nameItem)
                elif str(column[xyz])in("SITE0","SITE1","SITE2","SITE3","SITE4","SITE5","SITE6","SITE7","SITE8"):
                    nameItem = QTableWidgetItem()
                    nameItem.setData(Qt.DisplayRole,int(it['SITE'][str(column[xyz])[-1:]]['FAILS']))
                    w.item_su.setItem(row,xyz,nameItem)         
                else:
                    w.item_su.setItem(row,xyz,QTableWidgetItem(str(it[column[xyz]])))
            except:w.item_su.setItem(row,xyz,QTableWidgetItem("NoData"))
        only20000 = only20000 + 1
        try:
            if only20000>int(w.Item_num.text()):break#有輸入要列出幾個item就會執行這段code>>不然就參照預設(20000)
        except:
            if only20000>20000:break
    #w.item_su.horizontalHeader().setSectionResizeMode(QHeaderView.Stretch)
    #w.item_su.horizontalHeader().setSectionResizeMode(0, QHeaderView.Interactive)
    w.Status.setText('Status: 搜尋item完成')
    w.item_su.setSortingEnabled(True)#排列功能開啟
    w.Search_scatter.setEnabled(2)#搜尋完item後就可以開啟畫scatter功能
    w.pushButton_Group_By.setEnabled(2)
    w.pushButton_CID_Search_2.setEnabled(2)


def item_search():
    w.stackedWidget_su_it.setCurrentIndex(2)#將畫面轉到item
    while w.item_search_su.rowCount() > 0: w.item_search_su.removeRow(w.item_search_su.rowCount()-1) #先清空TABLE
    w.item_search_su.setSortingEnabled(False)#搜尋前>>先讓table不可排序>>因為可排序後搜尋會導致有些格子為空值
    if w.lineEdit_item.text()!='' and item_all!='':
        item_search_list,list_temp=[],[]
        list_temp=(str(w.lineEdit_item.text()).split(','))
        for ppap in list_temp:#如果有人key"240-247">>要變成[240,241,242...247]
            if ppap.find("-")!=-1:
                start,end= int(ppap[:ppap.find("-")]),int(ppap[ppap.find("-")+1:len(ppap)])
                while start <=end:
                    item_search_list.append(start)
                    start = start + 1
            else:item_search_list.append(ppap)
        #print(item_search_list)   
        for it_need in item_search_list:
            item = QTableWidgetItem()
            item.setFlags(item.flags())
            item.setCheckState(Qt.Checked)
            for it in item_all:
                if str(it['TEST_NUM']) == str(it_need):
                    row = w.item_search_su.rowCount()
                    w.item_search_su.setRowCount(row + 1)
                    w.item_search_su.setItem(row,0,item)
                    #column = ['cl','TEST_NUM','TEST_TXT','FAILS','EXECS','LTL','UTL',"SITE0","SITE1","SITE2","SITE3","SITE4","SITE5","SITE6","SITE7","SITE8"]
                    for xyz in range(1,len(column)):
                        try:
                            if str(column[xyz])=='FAILS' or str(column[xyz])=='TEST_NUM' or str(column[xyz])=='EXECS' or str(column[xyz])=='SEQ':
                                nameItem = QTableWidgetItem()
                                nameItem.setData(Qt.DisplayRole,int(it[column[xyz]]))
                                w.item_search_su.setItem(row,xyz,nameItem)
                            elif str(column[xyz])in("SITE0","SITE1","SITE2","SITE3","SITE4","SITE5","SITE6","SITE7","SITE8"):
                                nameItem = QTableWidgetItem()
                                nameItem.setData(Qt.DisplayRole,int(it['SITE'][str(column[xyz])[-1:]]['FAILS']))
                                w.item_search_su.setItem(row,xyz,nameItem)   
                            else:
                                w.item_search_su.setItem(row,xyz,QTableWidgetItem(str(it[column[xyz]])))
                        except:w.item_search_su.setItem(row,xyz,QTableWidgetItem("No Data"))
                    break
    w.item_search_su.setSortingEnabled(True)#排列功能開啟


def search_group_summary():
    global rs,AuthToken,item_all,findata,findata_tester,findata_id,findata_LB,findata_diff,findata_device,findata_TP_rev,findata_Handler,findata_time,httpm
    Group_by_what = int(str(w.comboBox_GroupBy.currentText())[0])#要by甚麼
    w.stackedWidget_su_it.setCurrentIndex(1)#將畫面轉到item
    while w.table_group_summary.rowCount() > 0: w.table_group_summary.removeRow(w.table_group_summary.rowCount()-1) #先清空TABLE
    w.table_group_summary.setSortingEnabled(False)#搜尋前>>先讓table不可排序>>因為可排序後搜尋會導致有些格子為空值
    #================================搜尋要Group的item================================
    #先看是要畫那些item
    user_item=[]
    for row in range(w.item_su.rowCount()):
        if w.item_su.item(row,0).checkState()==2:
            for num in range(len(item_all)):
                if int(item_all[num]['TEST_NUM'])==int(w.item_su.item(row,1).text()):
                    user_item.append(num)
                    break
    for row in range(w.item_search_su.rowCount()):
        if w.item_search_su.item(row,0).checkState()==2:
            for num in range(len(item_all)):
                if int(item_all[num]['TEST_NUM'])==int(w.item_search_su.item(row,1).text()):
                    user_item.append(num)
                    break
    if len(user_item)==0:
        w.Status.setText('Status: 沒偵察到可畫item...')
        alert('沒偵察到可畫item...')
        return
    #================================將現有的summary整理成一個個的Group(JOB_REV/LOAD_ID/LOT_ID/NODE_NAM/HAND_ID/ENG_ID)================================
    findata_time_temp=[]#將時間 "2019-06-05T14:05:11" 變成 "2019-06-05"
    for tttemp in findata_time:findata_time_temp.append(tttemp[:10])
    Group_By_column = ["Tester","Diff","LB","lotid","TP_rev","Handler"]
    Group_By_column = Group_By_column[0+Group_by_what:1+Group_by_what]#用Group_by_what取要用哪個因子
    df_Group_By_list=pd.DataFrame({"file":findata,"Time":findata_time_temp,"Tester":findata_tester,"Diff":findata_diff,"LB":findata_LB,"lotid":findata_id,"TP_rev":findata_TP_rev,"Handler":findata_Handler})
    Group_By_list = [findata_tester,findata_id,findata_LB,findata_diff,findata_TP_rev,findata_Handler]
    Group_By = []#放 ex: [DX1,DX2,DX3,L04DIAXMT95TQ21000001,L04DIAXMT95TQ21000002,....]
    Group_File = []#放對應Group_By有哪些file,如陣列的第一格就放都是為DX1的file的陣列 ex: [[file1,file2],[file1,file3],...]
    Group_File_time = []#放對應Group_By中所有file中最新的日期 ex:[[201906271444],[201905091200],...]
    for GB_tilte in Group_By_column:
        Group_By = Group_By + list(df_Group_By_list[GB_tilte].unique())
        for GBy in list(df_Group_By_list[GB_tilte].unique()):
            Group_File.append(list(df_Group_By_list[df_Group_By_list[GB_tilte]==GBy]["file"]))
            Group_File_time.append(max(list(df_Group_By_list[df_Group_By_list[GB_tilte]==GBy]["Time"])))
        
    #================================開始一個個item搜尋其Group================================
    #選擇哪一場>>中華\銅鑼
    Factory = str(w.factory.currentText())
    if Factory=='中華': httpm = '******'
    else: httpm = '******'
    for userneed in user_item:
        for GB in range(len(Group_By)):
            time.sleep(1)
            w.Status.setText('loadGroup'+str(Group_By[GB])+'資訊中...')
            print('loadGroup'+str(Group_By[GB])+'資訊中...')
            itemname = str(item_all[userneed]['TEST_TXT'])
            itemnum = str(item_all[userneed]['TEST_NUM'])
            itemlowlimit = str(item_all[userneed]['LTL'])
            itemhighlimit = str(item_all[userneed]['UTL'])
            if itemlowlimit == 'None' : itemlowlimit='{}'
            if itemhighlimit == 'None' : itemhighlimit='{}'
            try_time = 0
            try_max_time = 5
            fuckyouhaha=0
            while try_time<try_max_time:#至少download資料要先成功>>才能有後續動作
                try:
                    time.sleep(0.5)
                    js_search3='''******'''
                    js4=json.loads(js_search3)
                    js4["files"]=Group_File[GB]
                    header4={'******'}
                    res4=rs.post(httpm+'/api/datalog/filter/test/accu',headers=header4,json=js4)
                    Group_datalog=res4.json()
                    fuckyouhaha = Group_datalog[0]['EXECS']
                    break
                except MemoryError:
                    w.Status.setText('Status: 昱冠回傳檔案太大..請減少資料量or單一item畫畫看')
                    alert('昱冠回傳檔案太大..請減少資料量or單一item畫畫看')
                    try_time = try_max_time  
                except Exception:
                    try_time = try_time + 1
            if try_time>=try_max_time:#超過一定次數~~就放棄(有時程式或昱冠會卡住)
                w.Status.setText('Status: '+str(Group_By[GB])+'接收昱冠資料失敗..skip此Group..繼續下一個Group')
                alert(str(Group_By[GB])+'接收昱冠資料失敗..skip此Group..繼續下一個Group')
                continue
            #================================羅列item================================
            #參照用,真正的設定在最下面>>column_group_summary = ['Group','Time','TEST_TXT','TEST_NUM','SEQ','EXECS','FAILS','MEAN','CPK','MAX','MIN','LTL','UTL']
            for it in Group_datalog:
                row = w.table_group_summary.rowCount()
                w.table_group_summary.setRowCount(row + 1)
                w.table_group_summary.setItem(row,0,QTableWidgetItem(str(Group_By[GB])))       
                for xyz in range(1,len(column_group_summary)):
                    try:
                        if str(column_group_summary[xyz])=='FAILS' or str(column_group_summary[xyz])=='TEST_NUM' or str(column_group_summary[xyz])=='EXECS' or str(column_group_summary[xyz])=='SEQ':
                            nameItem = QTableWidgetItem()
                            nameItem.setData(Qt.DisplayRole,int(it[column_group_summary[xyz]]))
                            w.table_group_summary.setItem(row,xyz,nameItem)                    
                        elif str(column_group_summary[xyz])=='MEAN' or str(column_group_summary[xyz])=='CPK':
                            try:w.table_group_summary.setItem(row,xyz,QTableWidgetItem(str("%.2f"%it[column_group_summary[xyz]])))
                            except:w.table_group_summary.setItem(row,xyz,QTableWidgetItem("No Data"))       
                        elif str(column_group_summary[xyz])=='Time':
                            w.table_group_summary.setItem(row,xyz,QTableWidgetItem(str(Group_File_time[GB])))
                        else:
                            w.table_group_summary.setItem(row,xyz,QTableWidgetItem(str(it[column_group_summary[xyz]])))
                    except:w.table_group_summary.setItem(row,xyz,QTableWidgetItem("NoData"))

    #================================畫圖(但只有勾一個item才畫)================================
    if len(user_item)==1:
        df_tem=pd.DataFrame(columns=column_group_summary)
        for row in range(w.table_group_summary.rowCount()):
            temp_list=[]
            for col in range(w.table_group_summary.columnCount()):
                temp_list.append(w.table_group_summary.item(row,col).text())
            df_tem.loc[row]=temp_list
            df_tem['FAILS']=pd.to_numeric(df_tem['FAILS'])
            df_tem = df_tem.sort_values(by=['Time'],ascending=False)
            df_tem = df_tem.reset_index(drop=True)
        w.Mplwidget.canvas.axes.clear()
        #for coll in df_tem['Group'][df_tem['FAILS']>0].unique():
        #    w.Mplwidget.canvas.axes.bar(coll,df_tem['FAILS'][df_tem['Group']==coll])
        #x = np.arange(len(df_tem['Group'][df_tem['FAILS']>0].unique()))
        for coll in df_tem['Group']:
            w.Mplwidget.canvas.axes.bar(coll,df_tem['FAILS'][df_tem['Group']==coll])
        x = np.arange(len(df_tem['Group']))
        rotation_nu,x_size=15,8
        if len(x)<10:rotation_nu,x_size=20,6
        elif len(x)<20:rotation_nu,x_size=30,4
        w.Mplwidget.canvas.axes.tick_params(axis='y',labelsize=8)
        if len(x)<20:w.Mplwidget.canvas.axes.tick_params(axis='x',labelrotation=rotation_nu,labelsize=x_size)#項目個數小於20才顯示
        w.Mplwidget.canvas.draw()

    w.stackedWidget_su_it.setCurrentIndex(1)#將畫面轉到Group
    w.table_group_summary.setSortingEnabled(True)
    w.Status.setText('Status: 完成..')

def search_scatter():
    global rs,AuthToken,summary_need,item_all,findata,findata_romcod,findata_unit,findata_tester,findata_id,findata_time,findata_TP,findata_LB,findata_diff,findata_device,httpm
    w.stackedWidget_su_it.setCurrentIndex(2)#將畫面轉到item
    if item_all=='' or summary_need=='':
        w.Status.setText('無任何資料可畫圖')
        alert('無任何資料可畫圖')
        return
    #================================先確認初步的設定================================
    if w.By_Tester.isChecked():separate=0
    elif w.By_diff.isChecked():separate=1
    elif w.By_LB.isChecked():separate=2
    elif w.By_lot.isChecked():separate=3
    else:separate=5#當作勾勾>> 5>>就是都不分

    if w.CB_site0.isChecked()==False and w.CB_site1.isChecked()==False and w.CB_site2.isChecked()==False and w.CB_site3.isChecked()==False and w.CB_site4.isChecked()==False and w.CB_site5.isChecked()==False and w.CB_site6.isChecked()==False and w.CB_site7.isChecked()==False and w.CB_site8.isChecked()==False:
        alert('請~~至少要勾選一個site')#至少site1-8要勾選一個
        return
    if w.checkBox_show.isChecked():show_yes = False #不存ppt但是否show圖
    else:show_yes = True
    if w.checkBox_PPT.isChecked():output_ppt_yes = True #是否輸出ppt
    else:output_ppt_yes = False
    if w.Need_Combine.isChecked():separate_scatter_combine = True #是否將不同Tester(LB/Diff/Lotid)畫在同一張圖
    else:separate_scatter_combine = False
    if w.Need_Histogram.isChecked():histogram_need = True #是否畫histogram
    else:histogram_need = False
    #if w.All_in_One.isChecked():allinone = True #是否畫allinone
    #else:allinone = False
    
    #================================要by甚麼分,tester/diffusion/LB/program================================
    if separate!=5:
        w.Status.setText('載入資料中..請稍後...')
        print('載入資料中..請稍後...')
        separate_by=['NODE_NAM','ENG_ID','LOAD_ID','LOT_ID']
        separate_by_name=['']#放by xxx分之後,每個分類的名稱(如by Tester>>那就放:DX09,DX-010,DX-011)
        #需要長名/短名/json檔(所以勢必還要再重新搜尋一次)
        analysis_sep = pd.DataFrame({"File_name":findata,"LOT_ID":findata_id,"NODE_NAM":findata_tester,"ENG_ID":findata_diff,"LOAD_ID":findata_LB,"JOB_NAM":findata_TP,"Units":findata_unit})
        #存放by各種tester/diffusion/LB/program的 summary
        if len(analysis_sep[str(separate_by[separate])].unique().tolist())>1:#至少機台要超過2台才有分的意義
            x2=1
            for x1 in analysis_sep[str(separate_by[separate])].unique():
                separate_by_name.append(x1)
                locals()['separate_file_%s'%x2]=[]
                locals()['separate_file_%s'%x2]=list(analysis_sep[analysis_sep[str(separate_by[separate])]==x1].File_name)#放長名
                locals()['separate_id_%s'%x2]=[]
                locals()['separate_id_%s'%x2]=list(analysis_sep[analysis_sep[str(separate_by[separate])]==x1].LOT_ID)#放短名
                locals()['separate_summary_%s'%x2] = []
                for xyz in summary_need:
                    if xyz[str(separate_by[separate])]==x1:
                        locals()['separate_summary_%s'%x2].append(xyz)#放summary
                x2 = x2 + 1
        else:
            alert('你只有 '+str(analysis_sep[str(separate_by[separate])].unique().tolist())+' 一個!!\n學別人by '+str(separate_by[separate]))
            return
        #==將separate_file_x 與 separate_id_x 按照(xxxxxxxx_20180814203623.std.gz)的時間去排列如果filename的格式為如此的話
        for ttemp in range(1,x2):
            #print(locals()['separate_file_%s'%ttemp])
            #print(locals()['separate_id_%s'%ttemp])
            #print('============================')
            try:#如果filename的格式不符合的話>>error
                RIT_list=[]#RIT>>Right_in_Time
                for RIT in locals()['separate_file_%s'%ttemp]:
                    RIT_list.append(str(int(RIT[RIT.find('.std')-14:RIT.find('.std')])))
                RIT_temp_df=pd.DataFrame({"File_name":locals()['separate_file_%s'%ttemp],"LOT_ID":locals()['separate_id_%s'%ttemp],"RIT":RIT_list})
                RIT_temp_df=RIT_temp_df.sort_values(by=['RIT'],ascending=True)#照時間排序
                locals()['separate_file_%s'%ttemp]=list(RIT_temp_df.File_name)#把排序好時間的File_name放回去
                locals()['separate_id_%s'%ttemp]=list(RIT_temp_df.LOT_ID)#把排序好時間的LOT_ID放回去
                #print(locals()['separate_file_%s'%ttemp])
                #print(locals()['separate_id_%s'%ttemp])
                #print('============================')
            except:print(str(locals()['separate_file_%s'%ttemp])+'\n格式不符')
                    
    #===存放all 的summary
    separate_file_0=findata
    separate_id_0=findata_id
    separate_summary_0=summary_need
    item_0=item_all
    #===作時間排序>>將separate_file_0 與 separate_id_0 按照(xxxxxxxx_20180814203623.std.gz)的時間去排列如果filename的格式為如此的話
    #print(separate_file_0)
    #print(separate_id_0)
    #print('============================')
    try:#如果filename的格式不符合的話>>error
        RIT_list=[]#RIT>>Right_in_Time
        for RIT in separate_file_0:
            #print(str(int(RIT[RIT.find('.std')-14:RIT.find('.std')])))
            RIT_list.append(str(int(RIT[RIT.find('.std')-14:RIT.find('.std')])))
        RIT_temp_df=pd.DataFrame({"File_name":separate_file_0,"LOT_ID":separate_id_0,"RIT":RIT_list})
        RIT_temp_df=RIT_temp_df.sort_values(by=['RIT'],ascending=True)#照時間排序
        separate_file_0=list(RIT_temp_df.File_name)#把排序好時間的File_name放回去
        separate_id_0=list(RIT_temp_df.LOT_ID)#把排序好時間的LOT_ID放回去
        #print(separate_file_0)
        #print(separate_id_0)
        #print('============================')
    except:print(str(separate_file_0)+'\n讀取File_name裡面的時間~格式不符~無法讀取')

    #================================開始搜尋/畫item================================
    #先看是要畫那些item
    user_item=[]
    for row in range(w.item_su.rowCount()):
        if w.item_su.item(row,0).checkState()==2:
            for num in range(len(item_all)):
                if int(item_all[num]['TEST_NUM'])==int(w.item_su.item(row,1).text()):
                    user_item.append(num)
                    break
    for row in range(w.item_search_su.rowCount()):
        if w.item_search_su.item(row,0).checkState()==2:
            for num in range(len(item_all)):
                if int(item_all[num]['TEST_NUM'])==int(w.item_search_su.item(row,1).text()):
                    user_item.append(num)
                    break
    if len(user_item)==0:
        w.Status.setText('Status: 沒偵察到可畫item...')
        alert('沒偵察到可畫item...')
        return
    
    for userneed in user_item:
        time.sleep(1)
        w.Status.setText('load資訊中...')
        print('load資訊中...')
        #單批user要的item
        #寫入要搜尋的summary(json檔)
        itemname = str(item_all[userneed]['TEST_TXT'])
        itemnum = str(item_all[userneed]['TEST_NUM'])
        itemlowlimit = str(item_all[userneed]['LTL'])
        itemhighlimit = str(item_all[userneed]['UTL'])
        if itemlowlimit == 'None' : itemlowlimit='{}'
        if itemhighlimit == 'None' : itemhighlimit='{}'
        
        #放使用者是否有更改limit
        try:temp_LTL,temp_UTL=float(w.lineEdit_XMIN.text()),float(w.lineEdit_XMAX.text())
        except:temp_LTL,temp_UTL=0,0
        
        if separate!=5:#>>byxxx分的話
            if len(analysis_sep[str(separate_by[separate])].unique().tolist())>1:#至少機台要超過2台才有分的意義 且separate非5
                separate_num=len(list(analysis_sep[str(separate_by[separate])].unique()))+1#0>>for all / 1~xx>>看有幾台測試機
            else:
                alert('抱歉您只有'+str(analysis_sep[str(separate_by[separate])].unique().tolist()[0])+'\n請將By xxx的勾選清除')#有勾選byxxx卻有只有一個選項,則直接不允許繼續做下去
                return
                #separate_num=1
        else:separate_num=1

        #如果確定是要輸出ppt則直接不show圖>>而因為show出來的圖要用分圖,但存的圖要用單一圖>>所以才需要此參數
        if output_ppt_yes: save_yes=True
        else:save_yes=False
        while True:
            if save_yes:prs = Presentation()#存圖，才要在最前面定義這些>>不存圖，也就不需要ppt了
                
            #===================================================非combine再一起的===================================================
            if separate_scatter_combine==False:
                if save_yes==False:plt.figure(figsize=(14,7))
                for x2 in range(separate_num):
                    if save_yes==False and x2==0 and separate!=5: continue#如果不储存(只show圖意思)且是有勾選"By...."就不畫all(第一張)
                    if save_yes==False and x2>=5:break#如果不储存(只show圖意思)但又超過4張>>第4張之後就不畫
                    if save_yes==False:
                        if x2==1:plt.subplot(2,2,1)
                        elif x2==2:plt.subplot(2,2,2)
                        elif x2==3:plt.subplot(2,2,3)
                        elif x2==4:plt.subplot(2,2,4)
                        #subplots_adjust(left=None, bottom=None, right=None, top=None,wspace=0.2, hspace=0.2)
                    #選擇哪一場>>中華\銅鑼
                    Factory = str(w.factory.currentText())
                    if Factory=='中華': httpm = '******'
                    else: httpm = '******'
                    try_time=0
                    try_max_time=5#超過一定次數~~就放棄(有時程式或昱冠會卡住)
                    fuckyouasshole=0
                    while try_time<try_max_time:#至少download資料要先成功>>才能有後續動作
                        try:
                            time.sleep(0.5)
                            js_search3='''******'''
                            #"PAdata":true / "minMax":{"MIN":0,"MAX":578728.5}
                            js4=json.loads(js_search3)
                            js4['searchOpt']['files']=locals()['separate_file_%s'%x2]
                            js4['FILE_NAME']=locals()['separate_file_%s'%x2]
                            js4['searchOpt']['tests'][0]['TEST_TXT']=item_all[userneed]['TEST_TXT']#有些item Name裡面有特殊字元~所以必須這樣放進去
                            header4={'******'}
                            res4=rs.post(httpm+'/api/datalogAnalysis/PApageData',headers=header4,json=js4)
                            datalog=res4.json()
                            fuckyouasshole = datalog['inputData'][0]['DATA'][0]['HB']
                            break
                        except MemoryError:
                            w.Status.setText('Status: 昱冠回傳檔案太大..請減少資料量or單一item畫畫看')
                            alert('昱冠回傳檔案太大..請減少資料量or單一item畫畫看')
                            try_time = try_max_time  
                        except Exception:
                            try_time = try_time + 1
                    if try_time>=try_max_time:
                        w.Status.setText('Status: (item'+str(itemnum)+')(x2='+str(x2)+')接收昱冠資料失敗..繼續下一個item..or結束')
                        alert('(item'+str(itemnum)+')(x2='+str(x2)+')接收昱冠資料失敗..繼續下一個item..or結束')
                        try:del js_search3,js4,res4,datalog#獲取資料失敗>>就全部變數消除
                        except:pass
                        continue
                    #初始化
                    site0y,site1y,site2y,site3y,site4y,site5y,site6y,site7y,site8y=[],[],[],[],[],[],[],[],[]
                    site0x,site1x,site2x,site3x,site4x,site5x,site6x,site7x,site8x=[],[],[],[],[],[],[],[],[]
                    HLLL,HL,LL,sitey_all=[],[],[],[]
                    testtest,Max_value,Min_value=0.0,0.0,0.0#取得最大最小值
                    file_std,x_value_temp = '',0
                    file_std_name ,file_std_name_x ,file_std_name_note=[],[],[]
                    #load rowData
                    for RIT_num in locals()['separate_file_%s'%x2]:#對時間用的
                        for xyz in range(len(datalog['inputData'])):
                            if RIT_num==str(datalog['inputData'][xyz]['_id']['FILE_NAME']):#對時間用的
                                for ppap in range(len(datalog['inputData'][xyz]['DATA'])):
                                    try:
                                        data = datalog['inputData'][xyz]['DATA'][ppap]
                                        try:testtest = float(data['VAL'])
                                        except:continue#如果值是None之類的~就直接跳過
                                        if file_std != data['FILE_NAME']:#for 註解/summary間隔線用
                                            file_std = data['FILE_NAME']
                                            file_std_name.append(file_std)
                                            file_std_name_x.append(int(x_value_temp))
                                            file_std_name_note.append('')
                                        if data['SITE']==0:
                                            site0y.append(data['VAL'])
                                            site0x.append(int(x_value_temp))
                                        if data['SITE']==1:
                                            site1y.append(data['VAL'])
                                            site1x.append(int(x_value_temp))
                                        if data['SITE']==2:
                                            site2y.append(data['VAL'])
                                            site2x.append(int(x_value_temp))
                                        if data['SITE']==3:
                                            site3y.append(data['VAL'])
                                            site3x.append(int(x_value_temp))
                                        if data['SITE']==4:
                                            site4y.append(data['VAL'])
                                            site4x.append(int(x_value_temp))
                                        if data['SITE']==5:
                                            site5y.append(data['VAL'])
                                            site5x.append(int(x_value_temp))
                                        if data['SITE']==6:
                                            site6y.append(data['VAL'])
                                            site6x.append(int(x_value_temp))
                                        if data['SITE']==7:
                                            site7y.append(data['VAL'])
                                            site7x.append(int(x_value_temp))
                                        if data['SITE']==8:
                                            site8y.append(data['VAL'])
                                            site8x.append(int(x_value_temp))
                                        x_value_temp = x_value_temp + 1
                                    except:pass
                    sitey_all = site0y+site1y+site2y+site3y+site4y+site5y+site6y+site7y+site8y
                    Max_value = max(sitey_all)
                    Min_value = min(sitey_all)
                    MM=np.mean(sitey_all)#給標註的內容其y軸位置用
                    for xyz in range(len(datalog['inputData'])):
                        HL.append(datalog['inputData'][xyz]['HI_LIMIT'])
                        LL.append(datalog['inputData'][xyz]['LO_LIMIT'])
                        HLLL.append([datalog['inputData'][xyz]['LO_LIMIT'],datalog['inputData'][xyz]['HI_LIMIT']])
                    #print(MM)
                    abc=datetime.now()
                    abcpath='D:/python/'+str(abc)[0:4]+'-'+str(abc)[5:7]+'-'+str(abc)[8:10]+'-'+str(abc)[11:13]+str(abc)[14:16]+str(abc)[17:19]+'_T'+str(itemnum)+'_'+str(x2)#存圖用的流水號檔名                
                    if histogram_need==False:
                        #=============================================畫scatter===============================================
                        print('儲存scatter中...請稍後')
                        if save_yes:plt.figure(figsize=(14,7))#存>>>才是單張單張的畫
                        if len(site0y)>0 and w.CB_site0.isChecked():plt.scatter(site0x,site0y,s=15,c='m',marker='o',alpha=0.6,label='site0')
                        if len(site1y)>0 and w.CB_site1.isChecked():plt.scatter(site1x,site1y,s=15,c='g',marker='x',alpha=0.6,label='site1')
                        if len(site2y)>0 and w.CB_site2.isChecked():plt.scatter(site2x,site2y,s=15,c='r',marker='s',alpha=0.6,label='site2')
                        if len(site3y)>0 and w.CB_site3.isChecked():plt.scatter(site3x,site3y,s=15,c='blue',marker='^',alpha=0.6,label='site3')
                        if len(site4y)>0 and w.CB_site4.isChecked():plt.scatter(site4x,site4y,s=15,c='orange',marker='v',alpha=0.6,label='site4')
                        if len(site5y)>0 and w.CB_site5.isChecked():plt.scatter(site5x,site5y,s=15,c='c',marker='x',alpha=0.6,label='site5')
                        if len(site6y)>0 and w.CB_site6.isChecked():plt.scatter(site6x,site6y,s=15,c='chartreuse',marker='s',alpha=0.6,label='site6')
                        if len(site7y)>0 and w.CB_site7.isChecked():plt.scatter(site7x,site7y,s=15,c='chocolate',marker='o',alpha=0.6,label='site7')
                        if len(site8y)>0 and w.CB_site8.isChecked():plt.scatter(site8x,site8y,s=15,c='darkolivegreen',marker='v',alpha=0.6,label='site8')

                        #標上報表分界
                        for xyz in range(len(file_std_name_x)):
                            if int(xyz)!=0:
                                plt.axvline(file_std_name_x[xyz],color='magenta',linestyle='dashed',linewidth=0.5)
                            if temp_LTL==0 and temp_UTL==0:
                                try:plt.text((file_std_name_x[xyz+1]-file_std_name_x[xyz])/2+file_std_name_x[xyz],MM,file_std_name_note[xyz],fontsize=10,verticalalignment="top",horizontalalignment="right")
                                except:plt.text((x_value_temp-file_std_name_x[xyz])/2+file_std_name_x[xyz],MM,file_std_name_note[xyz],fontsize=10,verticalalignment="top",horizontalalignment="right")
                        #如果有更改limit時
                        if temp_LTL!=0 or temp_UTL!=0:
                            new_ticks = np.linspace(temp_LTL,temp_UTL, 10)
                            plt.yticks(new_ticks)
                            plt.ylim(temp_LTL,temp_UTL)
                            for xyz in range(len(file_std_name_x)):
                                try:plt.text((file_std_name_x[xyz+1]-file_std_name_x[xyz])/2+file_std_name_x[xyz],(temp_UTL-temp_LTL)/2+temp_LTL,file_std_name_note[xyz],fontsize=10,verticalalignment="top",horizontalalignment="right")
                                except:plt.text((x_value_temp-file_std_name_x[xyz])/2+file_std_name_x[xyz],(temp_UTL-temp_LTL)/2+temp_LTL,file_std_name_note[xyz],fontsize=10,verticalalignment="top",horizontalalignment="right")

                        #註解內容
                        if w.ONSEMI_note.isChecked():#如果有勾ONSEMI_報表標註
                            for ppap in range(len(file_std_name)):
                                for aapa in range(len(locals()['separate_file_%s'%x2])):
                                    if file_std_name[ppap][:len(file_std_name[ppap])-7]==locals()['separate_file_%s'%x2][aapa][:len(locals()['separate_file_%s'%x2][aapa])-7]:
                                        mtfk = str(locals()['separate_file_%s'%x2][aapa])
                                        file_std_name_note[ppap]=mtfk[mtfk.find('1-1')+4:mtfk.find('_',mtfk.find('1-1'))]+'\n'+mtfk[mtfk.find('_',mtfk.find('1-1'))+1:mtfk.find('_',mtfk.find('1-1'))+11]
                                        break
                        #標上註解
                        for xyz in range(len(file_std_name_x)):
                            if int(xyz)!=0:
                                plt.axvline(file_std_name_x[xyz],color='magenta',linestyle='dashed',linewidth=1)
                            if temp_LTL==0 and temp_UTL==0:
                                try:plt.text((file_std_name_x[xyz+1]-file_std_name_x[xyz])/2+file_std_name_x[xyz],MM,file_std_name_note[xyz],fontsize=10,verticalalignment="top",horizontalalignment="right")
                                except:plt.text((x_value_temp-file_std_name_x[xyz])/2+file_std_name_x[xyz],MM,file_std_name_note[xyz],fontsize=10,verticalalignment="top",horizontalalignment="right")


                        #標上HighLimit
                        for xyz in range(len(HL)):
                            if HL[xyz] != None:
                                #plt.axhline(HL[xyz],color='peru',linestyle = '-.',linewidth=1)
                                plt.axhline(HL[xyz],color='black',linestyle = '-.',linewidth=1)
                                plt.annotate(s='UTL = '+str(HL[xyz]),xy=(0,HL[xyz]),xytext=(0,HL[xyz]))
                        #標上LowLimit
                        for xyz in range(len(LL)):
                            if LL[xyz] != None:
                                plt.axhline(LL[xyz],color='black',linestyle = '-.',linewidth=1)
                                plt.annotate(s='LTL = '+str(LL[xyz]),xy=(0,LL[xyz]),xytext=(0,LL[xyz]))
                        plt.legend(loc='upper left',fontsize=8)#bbox_to_anchor=(0,0.3)                        
                        #標上TITLE
                        tttemp=[]
                        for yyyy in range(len(HLLL)):
                            if HLLL[yyyy]not in tttemp:
                                tttemp.append(HLLL[yyyy])
                        if save_yes==False and (x2==2 or x2==3 or x2==4):#1次SHOW4張圖的~~只讓第一張有Title
                            print('1次SHOW4張圖的~~只讓第一張有完整Title，其他張show分類即可')
                            plt.title(str(list(analysis_sep[str(separate_by[separate])].unique())[x2-1]))
                        else: 
                            if x2!=0:plt.title('Test#'+str(itemnum)+'::'+str(itemname)+'\nTest_Limits'+str(tttemp)+'-'+str(list(analysis_sep[str(separate_by[separate])].unique())[x2-1]))
                            else:plt.title('Test#'+str(itemnum)+'::'+str(itemname)+'\nTest_Limits'+str(tttemp))
                        plt.ylabel('Value')
                        plt.grid(True, linestyle = "-.", linewidth = 0.5)
                        plt.savefig(abcpath+'_scatter.png',dpi=300)

                    if histogram_need:
                        #=============================================畫histogram===============================================
                        if temp_LTL!=0 or temp_UTL!=0:histogram_range=[temp_LTL,temp_UTL]
                        else:histogram_range=[min(sitey_all),max(sitey_all)]
                        if save_yes:plt.figure(figsize=(14,7))#存>>>才是單張單張的畫
                        if len(site0y)>0 and w.CB_site0.isChecked():plt.hist(site0y,bins=100,rwidth = 0.75,alpha=0.6,range=histogram_range,color='m',label='site0')
                        if len(site1y)>0 and w.CB_site1.isChecked():plt.hist(site1y,bins=98 ,rwidth = 0.75,alpha=0.6,range=histogram_range,color='b',label='site1')
                        if len(site2y)>0 and w.CB_site2.isChecked():plt.hist(site2y,bins=96 ,rwidth = 0.75,alpha=0.6,range=histogram_range,color='g',label='site2')
                        if len(site3y)>0 and w.CB_site3.isChecked():plt.hist(site3y,bins=94 ,rwidth = 0.75,alpha=0.6,range=histogram_range,color='orange',label='site3')
                        if len(site4y)>0 and w.CB_site4.isChecked():plt.hist(site4y,bins=92 ,rwidth = 0.75,alpha=0.6,range=histogram_range,color='r',label='site4')
                        if len(site5y)>0 and w.CB_site5.isChecked():plt.hist(site5y,bins=90 ,rwidth = 0.75,alpha=0.6,range=histogram_range,color='c',label='site5')
                        if len(site6y)>0 and w.CB_site6.isChecked():plt.hist(site6y,bins=88 ,rwidth = 0.75,alpha=0.6,range=histogram_range,color='chartreuse',label='site6')
                        if len(site7y)>0 and w.CB_site7.isChecked():plt.hist(site7y,bins=86 ,rwidth = 0.75,alpha=0.6,range=histogram_range,color='chocolate',label='site7')
                        if len(site8y)>0 and w.CB_site8.isChecked():plt.hist(site8y,bins=84 ,rwidth = 0.75,alpha=0.6,range=histogram_range,color='pink',label='site8')
                        plt.ylabel('Count')

                        #標上HighLimit
                        for xyz in range(len(HL)):
                            if HL[xyz] != None:
                                plt.axvline(HL[xyz],color='black',linestyle = '-.',linewidth=1)
                                if HL[xyz]==LL[xyz]:plt.annotate(s='LTL = UTL = '+str(HL[xyz]),xy=(HL[xyz],0),xytext=(HL[xyz],0))
                                else:plt.annotate(s='UTL = '+str(HL[xyz]),xy=(HL[xyz],0),xytext=(HL[xyz],0))
                        #標上LowLimit
                        for xyz in range(len(LL)):
                            if LL[xyz] != None:
                                plt.axvline(LL[xyz],color='black',linestyle = '-.',linewidth=1)
                                if HL[xyz]!=LL[xyz]:plt.annotate(s='LTL = '+str(LL[xyz]),xy=(LL[xyz],0),xytext=(LL[xyz],0))

                        #如果有更改limit時
                        if temp_LTL!=0 or temp_UTL!=0:
                            new_ticks = np.linspace(round(temp_LTL,3),round(temp_UTL,3), 10)
                            #plt.xticks(new_ticks)
                            plt.xlim(temp_LTL,temp_UTL)
                        plt.legend(loc='upper left',fontsize=8)
                        #標上TITLE
                        tttemp=[]
                        for yyyy in range(len(HLLL)):
                            if HLLL[yyyy]not in tttemp:
                                tttemp.append(HLLL[yyyy])
                        if save_yes==False and (x2==2 or x2==3 or x2==4):#1次SHOW4張圖的~~只讓第一張有Title
                            print('1次SHOW4張圖的~~只讓第一張有完整Title，其他張show分類即可')
                            plt.title(str(list(analysis_sep[str(separate_by[separate])].unique())[x2-1]))
                        else: 
                            if x2!=0:plt.title('Test#'+str(itemnum)+'::'+str(itemname)+'\nTest_Limits'+str(tttemp)+'-'+str(list(analysis_sep[str(separate_by[separate])].unique())[x2-1]))
                            else:plt.title('Test#'+str(itemnum)+'::'+str(itemname)+'\nTest_Limits'+str(tttemp))
                        plt.savefig(abcpath+'_histogram.png',dpi=300)
                        
                    #如果是要存ppt(不show圖)輸出完圖後就清圖(減少暫存)
                    if save_yes:plt.close('all')
                    
                    #輸出ppt
                    if save_yes:
                        #選擇哪一場>>中華\銅鑼
                        Factory = str(w.factory.currentText())
                        if Factory=='中華': httpm = '******'
                        else: httpm = '******'
                        #得到site資訊
                        try_max_time=5
                        try_time = 0
                        fuckyouasshole=0
                        while try_time < try_max_time:
                            try:
                                js_search4='''******'''
                                js5=json.loads(js_search4)
                                js5['getDataOpt']['searchOpt']=js4['searchOpt']
                                js5['rawdata']=datalog['inputData']
                                for xyz in range(len(datalog['inputData'])):
                                    js5['rawdata'][xyz]['searchOpt']=js4['searchOpt']
                                js5['getDataOpt']['FILE_NAME']=locals()['separate_file_%s'%x2]
                                header5={'User-Authorization':AuthToken}
                                res5=rs.post(httpm+'/api/datalogAnalysis/bckPlot',headers=header5,json=js5)
                                sitedatalog=res5.json()
                                fuckyouasshole = sitedatalog['gridOpts']['rowData'][0]['Execs']
                                break
                            except MemoryError:
                                w.Status.setText('Status: 昱冠回傳檔案太大..請減少資料量or單一item畫畫看')
                                alert('昱冠回傳檔案太大..請減少資料量or單一item畫畫看')
                                try_time = try_max_time                          
                            except Exception:
                                try_time = try_time + 1
                        if try_time >=try_max_time :
                            w.Status.setText('Status: 做ppt所需的by site資訊..接收昱冠資料失敗..下一個or結束')
                            alert('做ppt所需的by site資訊..接收昱冠資料失敗..下一個or結束')
                            try:del js_search4,js5,res5,sitedatalog#獲取資料失敗>>就全部變數消除
                            except:pass
                            continue

                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        def ttv(slide,text,left,top,width,high,size,return_run):
                            shape1=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(high))#增加文字方框
                            #字型
                            run = shape1.text_frame.paragraphs[0].add_run()
                            run.text=text
                            font = run.font
                            font.name = 'Times New Roman'
                            font.size = Pt(size)
                            font.bold =True
                            font.color.rgb = RGBColor(0,0,255)
                            #框框內的樣式
                            fill = shape1.fill
                            fill.solid()
                            fill.fore_color.rgb = RGBColor(255, 255, 255)
                            fill.fore_color.brightness = 0
                            if return_run:
                                return run
                        text='Test#'+str(itemnum)+'  '+str(itemname)
                        if x2!=0:text=text+'-'+str(list(analysis_sep[str(separate_by[separate])].unique())[x2-1])
                        ttv(slide,text, 0 , 0.1 , 10 , 0.4 , 20,False)#字型
                        #------------讀site的excel------------
                        site_infor_column = ['Mean','Cp','Cpk','Execs','Fails','LowLimit','HighLimit']
                        df_site_info=pd.DataFrame(columns=['Site']+site_infor_column)
                        index = 0
                        for site_infor in sitedatalog['gridOpts']['rowData']:
                            for ch_x in range(9):
                                if int(str(site_infor['GroupField'])[-1:])==ch_x:#掃0-9有哪個SITE
                                    site_item_info=[]
                                    site_item_info.append(str(site_infor['GroupField']))
                                    for ch_y in range(2,9):#填入各SITE的各種值
                                        if str(site_infor_column[ch_y-2]) in ['Mean','Cp','Cpk']:
                                            try:site_item_info.append("%.2f" %site_infor[site_infor_column[ch_y-2]])
                                            except:site_item_info.append('NA')
                                        else:
                                            try:site_item_info.append(str(site_infor[site_infor_column[ch_y-2]]))
                                            except:site_item_info.append('NA')
                                    df_site_info.loc[index]=site_item_info
                                    index = index + 1
                        #------------添加site表格------------
                        rownum,colnum=5,8
                        tableTC = slide.shapes.add_table(rownum, colnum, Inches(0), Inches(5.2), Inches(10), Inches(0)).table
                        tableTC2 = slide.shapes.add_table(rownum, colnum, Inches(5), Inches(5.2), Inches(10), Inches(0)).table
                        for row in range(rownum):
                            for col in range(colnum):
                                cell=tableTC.rows[row].cells[col]
                                cell2=tableTC2.rows[row].cells[col]
                                paragraph=cell.text_frame.paragraphs[0]
                                paragraph2=cell2.text_frame.paragraphs[0]
                                paragraph.font.size = Pt(10)
                                paragraph2.font.size = Pt(10)
                                paragraph.font.bold = True
                                paragraph2.font.bold = True
                        for col in range(colnum):
                            tableTC.columns[col].width = Inches(0.58)
                            tableTC2.columns[col].width = Inches(0.58)
                            if col in [0,6,7]:
                                tableTC.columns[col].width = Inches(0.52)
                                tableTC2.columns[col].width = Inches(0.52)
                            if col in [1,2,3]:
                                tableTC.columns[col].width = Inches(0.75)
                                tableTC2.columns[col].width = Inches(0.75)
                        #先填入site table 的title
                        for col in range(1,colnum):
                            tableTC.cell(0,col).text=str(site_infor_column[col-1])
                            tableTC2.cell(0,col).text=str(site_infor_column[col-1])
                        for row in range(1,df_site_info.shape[0]+1):
                            for col in range(colnum):
                                if row<=4:tableTC.cell(row,col).text=df_site_info.iloc[row-1,col]
                                else:tableTC2.cell(row-4,col).text=df_site_info.iloc[row-1,col]
                        #------------添加圖片------------
                        if histogram_need==False:pic = slide.shapes.add_picture(abcpath+'_scatter.png', Inches(0), Inches(0.55), Inches(10), Inches(4.5))
                        elif histogram_need:pic = slide.shapes.add_picture(abcpath+'_histogram.png', Inches(0), Inches(0.55), Inches(10), Inches(4.5))
                    #清除記憶體
                    try:del site0y,site1y,site2y,site3y,site4y,site5y,site6y,site7y,site8y,site0x,site1x,site2x,site3x,site4x,site5x,site6x,site7x,site8x,HLLL,HL,LL,sitey_all,MM,Max_value,Min_value,file_std_name ,file_std_name_x ,file_std_name_note,file_std,js_search4,js5,res5,sitedatalog,js_search3,js4,res4,datalog
                    except:pass
            #===================================================畫combine再一起的===================================================
                    
            if separate_scatter_combine==True:
                if separate_num>18:#因為marker_list / color_list 只有用18個
                    separate_num=18
                    w.Status.setText('Status: 最多只畫18個,超過會error,不要問為什麼,就算是冠瑜也不給改')
                    alert('最多只畫18個,超過會error,不要問為什麼,就算是冠瑜也不給改')

                histogram_range=[0,0]#取得最大最小直>>for histogram用於range
                plt.figure(figsize=(14,7))#不管存或不存>>都只畫單張
                if separate_num==1:
                    w.Status.setText('Status: 請檢查是否您所勾選的"by..."有兩種以上 or 是否有勾選xxx')
                    print('請檢查是否您所勾選的"by..."有兩種以上 or 是否有勾選xxx')
                    plt.close('all') #error就清掉圖
                    break
                
                #取得最大最小直>>for histogram用於range
                if histogram_need:
                    #選擇哪一場>>中華\銅鑼
                    Factory = str(w.factory.currentText())
                    if Factory=='中華': httpm = '******'
                    else: httpm = '******'
                    try_max_time=5
                    try_time = 0
                    fuckyouhaha = 0
                    while try_time < try_max_time:
                        try:
                            rs=requests.session()
                            js_search2='******'
                            header3={'******'}
                            js3=json.loads(js_search2)
                            js3['files']=findata
                            res3=rs.post(httpm+'/api/datalog/filter/test/accu',headers=header3,json=js3)
                            fuckyouhaha = res3.json()[0]["FAILS"]#如果這邊就error就表示沒有抓到資料
                            break
                        except MemoryError:
                            w.Status.setText('Status: 昱冠回傳檔案太大..請減少資料量or單一item畫畫看')
                            alert('昱冠回傳檔案太大..請減少資料量or單一item畫畫看')
                            try_time = try_max_time  
                        except Exception:
                           try_time = try_time + 1
                           time.sleep(try_time)
                    if try_time >= try_max_time:
                        w.Status.setText('Status: 做histogram所需的最大最小值資訊..接收昱冠資料失敗')
                        alert('做histogram所需的最大最小值資訊..接收昱冠資料失敗')
                        try:
                            plt.close('all')
                            del js_search2,js3,res3#error就將變數清空
                        except:pass
                        return
                    for it in res3.json():
                        if itemname==str(it['TEST_TXT']):
                            histogram_range=[it['MIN'],it['MAX']]
                            #histogram_range_per=(it['MAX']-it['MIN'])/100#以免後面有些貨批的最大最小值超過這個範圍
                            print(str(it['TEST_TXT']))
                            print(histogram_range)
                            #print(histogram_range_per)
                            break                   

                #將上下限的標註放在for x2 in range(separate_num)之外>>for圖的title
                HLLL,tttemp=[],[]

                for x2 in range(separate_num):
                    if x2==0:continue #直接從x2==1開始(要by的第一個機台開始)>>x2是all的
                    #選擇哪一場>>中華\銅鑼
                    Factory = str(w.factory.currentText())
                    if Factory=='中華': httpm = '******'
                    else: httpm = '******'
                    try_time=0
                    try_max_time=5
                    fuckyouhaha = 0
                    while try_time<try_max_time:#至少download資料要先成功>>才能有後續動作
                        try:
                            time.sleep(0.5)                        
                            js_search3='''******'''
                            js4=json.loads(js_search3)
                            js4['searchOpt']['files']=locals()['separate_file_%s'%x2]
                            js4['FILE_NAME']=locals()['separate_file_%s'%x2]
                            js4['searchOpt']['tests'][0]['TEST_TXT']=item_all[userneed]['TEST_TXT']#有些item Name裡面有特殊字元~所以必須這樣放進去
                            header4={'User-Authorization':AuthToken}
                            res4=rs.post(httpm+'/api/datalogAnalysis/PApageData',headers=header4,json=js4)
                            datalog=res4.json()
                            fuckyouhaha = datalog['inputData'][0]['DATA'][0]['HB']
                            break
                        except MemoryError:
                            w.Status.setText('Status: 昱冠回傳檔案太大..請減少資料量or單一item畫畫看')
                            alert('昱冠回傳檔案太大..請減少資料量or單一item畫畫看')
                            try_time = try_max_time  
                        except Exception:
                            try_time = try_time + 1
                    if try_time>=try_max_time:#超過一定次數~~就放棄(有時程式或昱冠會卡住)
                        w.Status.setText('Status: 昱冠資料回傳失敗..(x2='+str(x2)+')'+str(analysis_sep[str(separate_by[separate])].unique().tolist()[0]))
                        alert('昱冠資料回傳失敗..(x2='+str(x2)+')'+str(analysis_sep[str(separate_by[separate])].unique().tolist()[0]))
                        try:
                            plt.close('all')
                            del js_search3,js4,res4,datalog#error就將變數清空
                        except:pass                        
                        return
                    #初始化
                    HL,LL,site_y,site_x,x_value_temp=[],[],[],[],0
                    MM = 0.0 #給標註的內容其y軸位置用
                    testtest,Max_value,Min_value=0.0,0.0,0.0#取得最大最小值

                    for RIT_num in locals()['separate_file_%s'%x2]:#對時間用的
                        for xyz in range(len(datalog['inputData'])):
                            if RIT_num==str(datalog['inputData'][xyz]['_id']['FILE_NAME']):#對時間用的
                                for ppap in range(len(datalog['inputData'][xyz]['DATA'])):
                                    try:
                                        data = datalog['inputData'][xyz]['DATA'][ppap]
                                        try:testtest = float(data['VAL'])
                                        except:continue#如果值是None之類的~就直接跳過
                                        if data['SITE']==0:
                                            site_y.append(data['VAL'])
                                            site_x.append(int(x_value_temp))
                                        if data['SITE']==1:
                                            site_y.append(data['VAL'])
                                            site_x.append(int(x_value_temp))
                                        if data['SITE']==2:
                                            site_y.append(data['VAL'])
                                            site_x.append(int(x_value_temp))
                                        if data['SITE']==3:
                                            site_y.append(data['VAL'])
                                            site_x.append(int(x_value_temp))
                                        if data['SITE']==4:
                                            site_y.append(data['VAL'])
                                            site_x.append(int(x_value_temp))
                                        if data['SITE']==5:
                                            site_y.append(data['VAL'])
                                            site_x.append(int(x_value_temp))
                                        if data['SITE']==6:
                                            site_y.append(data['VAL'])
                                            site_x.append(int(x_value_temp))
                                        if data['SITE']==7:
                                            site_y.append(data['VAL'])
                                            site_x.append(int(x_value_temp))
                                        if data['SITE']==8:
                                            site_y.append(data['VAL'])
                                            site_x.append(int(x_value_temp))
                                        x_value_temp = x_value_temp + 1
                                    except:pass
                    Max_value = max(site_y)
                    Min_value = min(site_y)
                    MMnocount=0
                    for xyz in range(len(datalog['inputData'])):
                        HL.append(datalog['inputData'][xyz]['HI_LIMIT'])
                        LL.append(datalog['inputData'][xyz]['LO_LIMIT'])
                        HLLL.append([datalog['inputData'][xyz]['LO_LIMIT'],datalog['inputData'][xyz]['HI_LIMIT']])

                    marker_list=['o','x','s','^','v','x','s','o','v','.','<','>','p','*','h','+','D','1']
                    color_list=['r','g','m','blue','orange','c','chartreuse','chocolate','pink','antiquewhite','aqua','darkmagenta','darkolivegreen','gold','honeydew','lightsalmon','plum','seashell']                        

                    if histogram_need==False:
                        #=============================================畫scatter===============================================
                        w.Status.setText('儲存scatter中...請稍後')
                        print('儲存scatter中...請稍後')
                        if len(site_y)>0:#有值才畫圖
                            plt.scatter(site_x,site_y,s=15,c=color_list[x2],marker=marker_list[x2],alpha=0.6,label=str(list(analysis_sep[str(separate_by[separate])].unique())[x2-1]))
                        #如果有更改limit時
                        if temp_LTL!=0 or temp_UTL!=0:
                            new_ticks = np.linspace(temp_LTL,temp_UTL, 10)
                            plt.yticks(new_ticks)
                            plt.ylim(temp_LTL,temp_UTL)
                        #標上HighLimit
                        for xyz in range(len(HL)):
                            if HL[xyz] != None:
                                plt.axhline(HL[xyz],color='black',linestyle = '-.',linewidth=1)
                                plt.annotate(s='UTL = '+str(HL[xyz]),xy=(0,HL[xyz]),xytext=(0,HL[xyz]))
                        #標上LowLimit
                        for xyz in range(len(LL)):
                            if LL[xyz] != None:
                                plt.axhline(LL[xyz],color='black',linestyle = '-.',linewidth=1)
                                plt.annotate(s='LTL = '+str(LL[xyz]),xy=(0,LL[xyz]),xytext=(0,LL[xyz]))
                        plt.legend(loc='upper left',fontsize=10)#bbox_to_anchor=(0,0.3)
                        #標上TITLE
                        for yyyy in range(len(HLLL)):
                            if HLLL[yyyy]not in tttemp:
                                tttemp.append(HLLL[yyyy])
                        plt.title('Test#'+str(itemnum)+'::'+str(itemname)+'\nTest_Limits'+str(tttemp))
                    if histogram_need:
                        #============================================畫histogram===============================================
                        print('儲存histogram中...請稍後')
                        w.Status.setText('儲存histogram中...請稍後')
                        if temp_LTL!=0 or temp_UTL!=0:histogram_range=[temp_LTL,temp_UTL]
                        if len(site_y)>0:#有值才畫圖
                            #if (temp_LTL==0 and temp_UTL==0) and (Max_value > histogram_range[1]):#如最大值筆range還大,則原本分100個間距,就要在多分出幾個
                            #    plt.hist(site_y,bins=(100-x2+int((Max_value-histogram_range[1])/histogram_range_per)),rwidth = 0.8,alpha=0.6,color=color_list[x2],range=[histogram_range[0],Max_value],label=str(list(analysis_sep[str(separate_by[separate])].unique())[x2-1]))
                            #elif (temp_LTL==0 and temp_UTL==0) and ( Min_value < histogram_range[0]):#同上述原因
                            #    plt.hist(site_y,bins=(100-x2+int((histogram_range[0]-Min_value)/histogram_range_per)),rwidth = 0.8,alpha=0.6,color=color_list[x2],range=[Min_value,histogram_range[1]],label=str(list(analysis_sep[str(separate_by[separate])].unique())[x2-1]))
                            #else:#都沒有就正常畫
                            plt.hist(site_y,bins=(100-x2),rwidth = 0.8,alpha=0.6,color=color_list[x2],range=histogram_range,label=str(list(analysis_sep[str(separate_by[separate])].unique())[x2-1]))
                                
                        #標上HighLimit
                        for xyz in range(len(HL)):
                            if HL[xyz] != None:
                                plt.axvline(HL[xyz],color='black',linestyle = '-.',linewidth=1)
                                if HL[xyz]==LL[xyz]:plt.annotate(s='LTL = UTL = '+str(HL[xyz]),xy=(HL[xyz],0),xytext=(HL[xyz],0))
                                else:plt.annotate(s='UTL = '+str(HL[xyz]),xy=(HL[xyz],0),xytext=(HL[xyz],0))
                        #標上LowLimit
                        for xyz in range(len(LL)):
                            if LL[xyz] != None:
                                plt.axvline(LL[xyz],color='black',linestyle = '-.',linewidth=1)
                                if HL[xyz]!=LL[xyz]:plt.annotate(s='LTL = '+str(LL[xyz]),xy=(LL[xyz],0),xytext=(LL[xyz],0))
                        #如果有更改limit時
                        if temp_LTL!=0 or temp_UTL!=0:
                            new_ticks = np.linspace(round(temp_LTL,3),round(temp_UTL,3), 10)
                            plt.xlim(temp_LTL,temp_UTL)
                        plt.legend(loc='upper left',fontsize=10)
                        #標上TITLE
                        for yyyy in range(len(HLLL)):
                            if HLLL[yyyy]not in tttemp:
                                tttemp.append(HLLL[yyyy])
                        plt.title('Test#'+str(itemnum)+'::'+str(itemname)+'\nTest_Limits'+str(tttemp))
                    #清除記憶體
                    try:del HL,LL,site_y,site_x,x_value_temp
                    except:pass
                abc=datetime.now()#存儲檔案時的流水編號
                abcpath='D:/python/'+str(abc)[0:4]+'-'+str(abc)[5:7]+'-'+str(abc)[8:10]+'-'+str(abc)[11:13]+str(abc)[14:16]+str(abc)[17:19]+'_T'+str(itemnum)+'_'+str(x2)#存圖用的流水號檔名
                if histogram_need==False:
                    plt.xlabel('variables x')
                    plt.ylabel('Value')
                    plt.grid(True, linestyle = "-.", linewidth = 0.5)
                    plt.savefig(abcpath+'_scatter.png',dpi=300)
                if histogram_need:
                    plt.xlabel('Value')
                    plt.ylabel('Count')
                    plt.savefig(abcpath+'_histogram.png',dpi=300)
                if save_yes:plt.close('all') #如果是要存ppt(不show圖)則存完圖就清圖(減少暫存)
                
                #輸出ppt
                if save_yes:            
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    def ttv(slide,text,left,top,width,high,size,return_run):
                        shape1=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(high))#增加文字方框
                        #字型
                        run = shape1.text_frame.paragraphs[0].add_run()
                        run.text=text
                        font = run.font
                        font.name = 'Times New Roman'
                        font.size = Pt(size)
                        font.bold =True
                        font.color.rgb = RGBColor(0,0,255)
                        #框框內的樣式
                        fill = shape1.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(255, 255, 255)
                        fill.fore_color.brightness = 0
                        if return_run:
                            return run
                    text='Test#'+str(itemnum)+'  '+str(itemname)
                    ttv(slide,text, 0 , 0.1 , 10 , 0.4 , 20,False)#字型
                    #------------讀site的excel------------
                    #得到site資訊
                    sep_infor_column = ['MEAN','CP','CPK','EXECS','FAILS','LTL','UTL']
                    df_sep_info=pd.DataFrame(columns=['_']+sep_infor_column)
                    index = 0
                    for x2 in range(1,separate_num):
                        #選擇哪一場>>中華\銅鑼
                        Factory = str(w.factory.currentText())
                        if Factory=='中華': httpm = 'http://kyeda.kyec.com.tw:3000'
                        else: httpm = 'http://kyeda-tl.kyec.com.tw:3000'
                        try_max_time=5
                        try_time = 0
                        while try_time < try_max_time:
                            try:
                                time.sleep(0.5) 
                                #得到site資訊
                                js5=js4['searchOpt']
                                js5['files']=locals()['separate_file_%s'%x2]
                                js5['statsType']='A'
                                header5={'User-Authorization':AuthToken}
                                res5=rs.post(httpm+'/api/datalog/filter/test/accu',headers=header5,json=js5)
                                fuckyouhaha = res5.json()[0]["FAILS"]
                                break
                            except MemoryError:
                                w.Status.setText('Status: 昱冠回傳檔案太大..請減少資料量or單一item畫畫看')
                                alert('昱冠回傳檔案太大..請減少資料量or單一item畫畫看>>昱冠不給力~幹嘛')
                                try_time = try_max_time  
                            except Exception:
                               try_time = try_time + 1
                               time.sleep(try_time)
                        if try_time >= try_max_time:
                            w.Status.setText('Status: 做ppt所需的bysite資訊..接收昱冠資料失敗')
                            alert('做ppt所需的bysite資訊..接收昱冠資料失敗>>你有嚐過失敗的滋味嗎?')
                            try:del js5,res5 #error就清空記憶體
                            except:pass
                            return
                        sitedatalog=res5.json()                            
                        sep_item_info=[]
                        sep_item_info.append(str(list(analysis_sep[str(separate_by[separate])].unique())[x2-1]))
                        for ch_x in range(7):
                            if str(sep_infor_column[ch_x]) in ['MEAN','CP','CPK']:
                                try:sep_item_info.append("%.2f" %sitedatalog[0][str(sep_infor_column[ch_x])])
                                except:sep_item_info.append('NA')
                            else:
                                try:sep_item_info.append(sitedatalog[0][str(sep_infor_column[ch_x])])
                                except:sep_item_info.append('NA')
                        df_sep_info.loc[index]=sep_item_info
                        index = index + 1
                    #------------添加site表格------------
                    rownum,colnum=5,8
                    tableTC = slide.shapes.add_table(rownum, colnum, Inches(0), Inches(5.2), Inches(10), Inches(0)).table
                    tableTC2 = slide.shapes.add_table(rownum, colnum, Inches(5), Inches(5.2), Inches(10), Inches(0)).table
                    for row in range(rownum):
                        for col in range(colnum):
                            cell=tableTC.rows[row].cells[col]
                            cell2=tableTC2.rows[row].cells[col]
                            paragraph=cell.text_frame.paragraphs[0]
                            paragraph2=cell2.text_frame.paragraphs[0]
                            paragraph.font.size = Pt(10)
                            paragraph2.font.size = Pt(10)
                            paragraph.font.bold = True
                            paragraph2.font.bold = True
                    for col in range(colnum):
                        tableTC.columns[col].width = Inches(0.58)
                        tableTC2.columns[col].width = Inches(0.58)
                        if col in [0,6,7]:
                            tableTC.columns[col].width = Inches(0.52)
                            tableTC2.columns[col].width = Inches(0.52)
                        if col in [1,2,3]:
                            tableTC.columns[col].width = Inches(0.75)
                            tableTC2.columns[col].width = Inches(0.75)
                    #先填入site table 的title
                    for col in range(1,colnum):
                        tableTC.cell(0,col).text=str(sep_infor_column[col-1])
                        tableTC2.cell(0,col).text=str(sep_infor_column[col-1])
                    try:#因為群組可能超過8個~~但ppt的table資訊只能放8個>>超過的~~拎杯沒爽show啦
                        for row in range(1,df_sep_info.shape[0]+1):
                            for col in range(colnum):
                                if row<=4:tableTC.cell(row,col).text=str(df_sep_info.iloc[row-1,col])[:10]
                                else:tableTC2.cell(row-4,col).text=str(df_sep_info.iloc[row-1,col])[:10]
                    except:pass
                    #------------添加圖片------------
                    if histogram_need==False:pic = slide.shapes.add_picture(abcpath+'_scatter.png', Inches(0), Inches(0.55), Inches(10), Inches(4.5))
                    if histogram_need:pic = slide.shapes.add_picture(abcpath+'_histogram.png', Inches(0), Inches(0.55), Inches(10), Inches(4.5))
                    df_sep_info.to_excel (abcpath+'_test.xlsx', index = None, header=True)
            if save_yes:
                try_it=1
                while os.path.isfile(abcpath+'_Result.pptx')!=True:#因為很常沒存ppt
                    prs.save(abcpath+'_Result.pptx')
                    time.sleep(try_it)
                    try_it = try_it + 1
                    if try_it>5:
                        w.Status.setText('Status: 實在抱歉~ppt依然沒儲存..')
                        alert('實在抱歉~ppt依然沒儲存..也可能是我不想幫你存')
                        break

            if save_yes==False:
                if show_yes:plt.show()
                plt.close('all')
                if output_ppt_yes:
                    save_yes=True
                    continue
            break
    w.Status.setText('Status: 完成..')

def Binning_Analysis():
    global rs,AuthToken,summary_need,item_all,httpm
    w.stackedWidget_su_it.setCurrentIndex(0)#將畫面轉到summary
    
    #確認是否有報表與是否有勾選
    if w.summary.rowCount()==0:
        w.Status.setText('Status: 無任何報表')
        alert('沒有任何報表~~你按個毛??')
        #return
    findata_need,findata_id_need,findata_tester_need,findata_diff_need,findata_LB_need=[],[],[],[],[]
    for aaa in range(w.summary.rowCount()):
        if w.summary.item(aaa,0).checkState()==2:
            findata_need.append(w.summary.item(aaa,column_summary.index('Filename')).text())
            findata_id_need.append(w.summary.item(aaa,column_summary.index('Lot_Id')).text())
            findata_tester_need.append(w.summary.item(aaa,column_summary.index('Tester')).text())
            findata_diff_need.append(w.summary.item(aaa,column_summary.index('Fab_Lot_id')).text())
            findata_LB_need.append(w.summary.item(aaa,column_summary.index('Load_Id')).text())
    if len(findata_need)==0:
        w.Status.setText('Status: 無任何報表被勾選..')
        alert('無任何報表被勾選..你都沒再看的嗎?哀~~~~')
        #return

    #要By甚麼畫
    By_what=0#0>>都不要 1>>Tester 2>>diff 3>>LB
    By_what=int(str(w.comboBox_Byxx.currentText())[0])

    #存圖用的流水號檔名
    abc=datetime.now()
    abcpath='D:/python/'+str(abc)[0:4]+'-'+str(abc)[5:7]+'-'+str(abc)[8:10]+'-'+str(abc)[11:13]+str(abc)[14:16]+str(abc)[17:19]

    #----------------------------------------------Bin Table 分析--------------------------------------------------
    w.Status.setText('Status: Bin Table 分析中...')
    #選擇哪一場>>中華\銅鑼
    Factory = str(w.factory.currentText())
    if Factory=='中華': httpm = 'http://kyeda.kyec.com.tw:3000'
    else: httpm = 'http://kyeda-tl.kyec.com.tw:3000'
    try_max_time=3
    try_time = 0
    while try_time<try_max_time:
        try:
            js_search_Bintable='******'
            header_Bintable={'******'}
            js_Bintable=json.loads(js_search_Bintable)
            js_Bintable['FILE_NAME']=findata_need
            res_Bintable=rs.post(httpm+'/api/binning/getData',headers=header_Bintable,json=js_Bintable)
            fuckyouhaha = res_Bintable.json()[0]["FAILS"]#如果這邊ERROR就表示沒抓到
            break
        except:
            try_time = try_time +1
            time.sleep(try_time)
    if try_time >=4:
        alert('昱冠Download失敗...')
        #return
    Bintable = res_Bintable.json()
    BAA = pd.DataFrame(columns=("FILE_NAME","LOTID","NODE_NAM","ENG_ID","LOAD_ID","BIN_CNT","BIN_NAM","BIN_NUM","SITE_NUM"))
    BTAindex = 0
    for abc in Bintable :
        for xyz in range(len(findata_need)):
            if abc['FILE_NAME']==findata_need[xyz]:
                for abc_d in abc['SITE_BIN']:
                    BAA.loc[BTAindex] = [findata_need[xyz],findata_id_need[xyz],findata_tester_need[xyz],findata_diff_need[xyz],findata_LB_need[xyz],abc_d["BIN_CNT"],abc_d["BIN_NAM"],abc_d["BIN_NUM"],abc_d["SITE_NUM"]]
                    BTAindex = BTAindex + 1
                break
        
    #輸出原始資料
    #BAA.to_excel(abcpath+'_Summary_bintable_oringal.xlsx', index = None, header=True)

    #如果要By Tester,但卻只有一台Tester,則中斷不做
    if (By_what==1 and len(BAA.NODE_NAM.unique())<=1) or (By_what==2 and len(BAA.ENG_ID.unique())<=1) or (By_what==3 and len(BAA.LOAD_ID.unique())<=1):
        if By_what==1:alert('抱歉只有一個'+str(BAA.NODE_NAM.unique().tolist()[0])+'\n請重新選擇要By??')
        if By_what==2:alert('抱歉只有一個'+str(BAA.ENG_ID.unique().tolist()[0])+'\n請重新選擇要By??')
        if By_what==3:alert('抱歉只有一個'+str(BAA.LOAD_ID.unique().tolist()[0])+'\n請重新選擇要By??')
        return

    #-----------轉換成每個sb的資訊(每個Site的fail數量)-----------
    column=["BIN_NUM","BIN_NAM","BIN_NU_NA","ELEMENT","FAILS"]
    Bintable_analysis_EDA = pd.DataFrame(columns=column)#將BAA轉化成Bintable_analysis_EDA(EDA上的Bintable格式)
    site_count = BAA.sort_values(by='SITE_NUM')['SITE_NUM'].unique()#總共有幾個site排序後放進去

    #為了取得每個sb的Fail數量(因為必須先排序)
    for_sb_count = pd.DataFrame(columns=["BIN_NUM","BIN_NAM","FAILS"])
    zzz_index = 0
    for zzz in BAA['BIN_NUM'].unique().tolist():
        for_sb_count.loc[zzz_index]=[int(zzz),str(BAA[BAA['BIN_NUM']==zzz]['BIN_NAM'].tolist()[0]),int(BAA[BAA['BIN_NUM']==zzz]['BIN_CNT'].sum())]
        zzz_index = zzz_index + 1
    for_sb_count = for_sb_count.sort_values(by=['FAILS'],ascending=False)
    for_sb_count = for_sb_count.reset_index(drop=True)
    sb_count = for_sb_count.BIN_NUM.astype('int') #總共有幾個sb(加上強制轉換成int)
    for abc in range(len(site_count)):#有多少個SITE,Bintable_analysis_EDA col 就多幾個(site fail數量 與 %)
        Bintable_analysis_EDA.insert(len(column)+abc*2,'SITE'+str(site_count[abc]),[])
        Bintable_analysis_EDA.insert(len(column)+abc*2+1,'SITE'+str(site_count[abc])+'_per',[])

    #彙整成excel並輸出
    Tester_item=''
    LB_item=''
    Diff_item=''
    BaE_index = 0
    for abc in range(len(sb_count)):#將資料填入Bintable_analysis_EDA
        #放入總資訊
        sb_count_site = []
        sb_count_site.append(sb_count[abc]) #放入BIN_NUM
        sb_count_site.append(BAA[BAA['BIN_NUM']==sb_count[abc]]['BIN_NAM'].unique()[0]) #放入BIN_NAM
        #sb_count_site.append(str(sb_count_site[0])+'_'+str(sb_count_site[1])) #放入BIN_NU_NA>>畫圖用的(顯示在y軸的文字)
        sb_count_site.append(str(sb_count_site[0])) #放入BIN_NU_NA>>畫圖用的(顯示在y軸的文字)
        sb_count_site.append('all')#總資訊並沒有by xxx 分>>所以放入all
        sb_count_site.append(BAA[BAA['BIN_NUM']==sb_count[abc]]['BIN_CNT'].sum())#放入某個SBIN總FAIL數量
        for xyz in range(len(site_count)): #放入各site的sb加總
            site_sum = BAA[(BAA['BIN_NUM']==sb_count[abc])&(BAA['SITE_NUM']==site_count[xyz])]['BIN_CNT'].sum()
            sb_count_site.append(site_sum)
            sb_count_site.append(round(float(site_sum/BAA['BIN_CNT'].sum()),5))
        Bintable_analysis_EDA.loc[BaE_index]=[sb_count_site[xx] for xx in range(len(site_count)*2+len(column))] #把資料放入Bintable_analysis_EDA
        BaE_index = BaE_index + 1
        
        #放入by LB/Diff/Tester資訊
        if len(BAA.LOAD_ID.unique())>1:
            for LB_item in BAA.LOAD_ID.unique():
                sb_count_site = []
                sb_count_site.append(sb_count[abc]) #放入BIN_NUM
                sb_count_site.append(BAA[BAA['BIN_NUM']==sb_count[abc]]['BIN_NAM'].unique()[0]) #放入BIN_NAM
                #sb_count_site.append(str(sb_count_site[0])+'_'+str(sb_count_site[1])) #放入BIN_NU_NA>>畫圖用的(顯示在y軸的文字)
                sb_count_site.append(str(sb_count_site[0])) #放入BIN_NU_NA>>畫圖用的(顯示在y軸的文字)
                sb_count_site.append(LB_item)
                sb_count_site.append(BAA[(BAA['BIN_NUM']==sb_count[abc])&(BAA['LOAD_ID']==LB_item)]['BIN_CNT'].sum())#放入某個SBIN總FAIL數量
                for xyz in range(len(site_count)): #放入各site的sb加總
                    site_sum = BAA[(BAA['BIN_NUM']==sb_count[abc])&(BAA['LOAD_ID']==LB_item)&(BAA['SITE_NUM']==site_count[xyz])]['BIN_CNT'].sum()
                    sb_count_site.append(site_sum)
                    sb_count_site.append(round(float(site_sum/BAA[(BAA['LOAD_ID']==LB_item)]['BIN_CNT'].sum()),5))
                Bintable_analysis_EDA.loc[BaE_index]=[sb_count_site[xx] for xx in range(len(site_count)*2+len(column))] #把資料放入Bintable_analysis_EDA
                BaE_index = BaE_index + 1
                
        if len(BAA.ENG_ID.unique())>1:
            for Diff_item in BAA.ENG_ID.unique():
                sb_count_site = []
                sb_count_site.append(sb_count[abc]) #放入BIN_NUM
                sb_count_site.append(BAA[BAA['BIN_NUM']==sb_count[abc]]['BIN_NAM'].unique()[0]) #放入BIN_NAM
                #sb_count_site.append(str(sb_count_site[0])+'_'+str(sb_count_site[1])) #放入BIN_NU_NA>>畫圖用的(顯示在y軸的文字)
                sb_count_site.append(str(sb_count_site[0])) #放入BIN_NU_NA>>畫圖用的(顯示在y軸的文字)
                sb_count_site.append(Diff_item)
                sb_count_site.append(BAA[(BAA['BIN_NUM']==sb_count[abc])&(BAA['ENG_ID']==Diff_item)]['BIN_CNT'].sum())#放入某個SBIN總FAIL數量
                for xyz in range(len(site_count)): #放入各site的sb加總
                    site_sum = BAA[(BAA['BIN_NUM']==sb_count[abc])&(BAA['ENG_ID']==Diff_item)&(BAA['SITE_NUM']==site_count[xyz])]['BIN_CNT'].sum()
                    sb_count_site.append(site_sum)
                    sb_count_site.append(round(float(site_sum/BAA[(BAA['ENG_ID']==Diff_item)]['BIN_CNT'].sum()),5))
                Bintable_analysis_EDA.loc[BaE_index]=[sb_count_site[xx] for xx in range(len(site_count)*2+len(column))] #把資料放入Bintable_analysis_EDA
                BaE_index = BaE_index + 1
                
        if len(BAA.NODE_NAM.unique())>1:
            for Tester_item in BAA.NODE_NAM.unique():
                sb_count_site = []
                sb_count_site.append(sb_count[abc]) #放入BIN_NUM
                sb_count_site.append(BAA[BAA['BIN_NUM']==sb_count[abc]]['BIN_NAM'].unique()[0]) #放入BIN_NAM
                #sb_count_site.append(str(sb_count_site[0])+'_'+str(sb_count_site[1])) #放入BIN_NU_NA>>畫圖用的(顯示在y軸的文字)
                sb_count_site.append(str(sb_count_site[0])) #放入BIN_NU_NA>>畫圖用的(顯示在y軸的文字)
                sb_count_site.append(Tester_item)
                sb_count_site.append(BAA[(BAA['BIN_NUM']==sb_count[abc])&(BAA['NODE_NAM']==Tester_item)]['BIN_CNT'].sum())#放入某個SBIN總FAIL數量
                for xyz in range(len(site_count)): #放入各site的sb加總
                    site_sum = BAA[(BAA['BIN_NUM']==sb_count[abc])&(BAA['NODE_NAM']==Tester_item)&(BAA['SITE_NUM']==site_count[xyz])]['BIN_CNT'].sum()
                    sb_count_site.append(site_sum)
                    sb_count_site.append(round(float(site_sum/BAA[(BAA['NODE_NAM']==Tester_item)]['BIN_CNT'].sum()),5))
                Bintable_analysis_EDA.loc[BaE_index]=[sb_count_site[xx] for xx in range(len(site_count)*2+len(column))] #把資料放入Bintable_analysis_EDA
                BaE_index = BaE_index + 1
    #Bintable輸出excel
    Bintable_analysis_EDA.to_excel (abcpath+'_Summary_bintable_all.xlsx', index = None, header=True)

    #---------------輸出圖-----------------
    def to_percent(temp, position):return '%1.2f'%(100*temp) + '%'#轉換y軸為%顯示
    pat = ['-', '+', 'x', '\\', '*', 'o', 'O', '.']#hack
    #['r','g','m','blue','orange','c','chartreuse','chocolate','pink']
    coll = ['m','green','red','blue','orange','yellow','aqua','chocolate','grey','aquamarine','cornsilk','sienna','fuchsia','seagreen','powderblue','aqua','purple','beige','moccasin']*10
    Bintable_histogram = Bintable_analysis_EDA[(Bintable_analysis_EDA['ELEMENT']=='all')&(Bintable_analysis_EDA['BIN_NUM']!=230)].head(11).reset_index(drop=True)
    if len(Bintable_histogram['BIN_NUM'].tolist()) <=10 : x_axis_sb_count=len(Bintable_histogram['BIN_NUM'].tolist())-1# -1  是因為 sb_count 裡面有包含sb1
    else : x_axis_sb_count = 10
    x = np.arange(x_axis_sb_count)
    #輸出總合圖
    if By_what==0:
        while os.path.isfile(abcpath+'_Bintable_analysis_all.png')!=True:
            y_bottom = np.array([0]*x_axis_sb_count)
            plt.figure(figsize=(14,7))
            for sisi in site_count:
                locals()['y%s'%sisi] = Bintable_histogram['SITE'+str(sisi)+'_per'].iloc[1:x_axis_sb_count+1]
                if y_bottom.sum()==0: plt.bar(x,locals()['y%s'%sisi],label = 'stie'+str(sisi),color=coll[sisi],width=0.8,align='center')
                else : plt.bar(x,locals()['y%s'%sisi],bottom=y_bottom,label = 'stie'+str(sisi),color=coll[sisi],width=0.8,align='center')
                y_bottom = y_bottom + np.array( locals()['y%s'%sisi])
            plt.gca().yaxis.set_major_formatter(FuncFormatter(to_percent))
            plt.yticks(fontsize=10)
            plt.xticks(x,np.array(Bintable_histogram["BIN_NU_NA"].iloc[1:x_axis_sb_count+1]),rotation=15,fontsize=7)
            plt.grid(True, linestyle = "-.", linewidth = 0.5)
            plt.legend(loc="upper right")
            plt.savefig(abcpath+'_Bintable_analysis_all.png',dpi=300)
            plt.show()
            plt.close('all')
        

    #輸出by LB資料
    if By_what==3 and len(BAA.LOAD_ID.unique())>1:
        while os.path.isfile(abcpath+'_Bintable_analysis_LB.png')!=True:
            plt.figure(figsize=(14,7))
            total_width, n = 0.8, len(BAA.LOAD_ID.unique())
            width = total_width / n
            wi = 0#x軸位置
            for LB_item in BAA.LOAD_ID.unique():
                y_bottom = np.array([0]*x_axis_sb_count)
                for sisi in site_count:
                    if wi==0: label =  'stie'+str(sisi)
                    else:label=None
                    locals()['y%s'%sisi] = Bintable_analysis_EDA[(Bintable_analysis_EDA['ELEMENT']==LB_item)]['SITE'+str(sisi)+'_per'].iloc[1:x_axis_sb_count+1]
                    if y_bottom.sum()==0: plt.bar(x+(wi-1)*width,locals()['y%s'%sisi],label = label,color=coll[sisi],width=width*0.9,align='center')
                    else : plt.bar(x+(wi-1)*width,locals()['y%s'%sisi],bottom=y_bottom,label = label,color=coll[sisi],width=width*0.9,align='center')
                    y_bottom = y_bottom + np.array( locals()['y%s'%sisi])
                for coor in range(x_axis_sb_count): plt.annotate('', xy=(x[coor]+(wi-1)*width,y_bottom[coor]),xytext=(x[coor]+(wi-1)*width,y_bottom[coor]+0.0001),arrowprops=dict(facecolor=coll[wi], shrink=0.5))
                plt.scatter(0,0,marker='v',label = LB_item,color=coll[wi])
                wi = wi + 1
            plt.gca().yaxis.set_major_formatter(FuncFormatter(to_percent))
            plt.yticks(fontsize=10)
            plt.xticks(x,np.array(Bintable_histogram["BIN_NU_NA"].iloc[1:x_axis_sb_count+1]),rotation=15,fontsize=7)
            plt.grid(True, linestyle = "-.", linewidth = 0.5)
            plt.legend(loc="upper right")
            plt.savefig(abcpath+'_Bintable_analysis_LB.png',dpi=300)   
            plt.show()
            plt.close('all')

    #輸出by Diff資料
    if By_what==2 and len(BAA.ENG_ID.unique())>1:
        while os.path.isfile(abcpath+'_Bintable_analysis_Diff.png')!=True:
            plt.figure(figsize=(14,7))
            total_width, n = 0.8, len(BAA.ENG_ID.unique())
            width = total_width / n
            wi = 0#x軸位置
            for Diff_item in BAA.ENG_ID.unique():
                y_bottom = np.array([0]*x_axis_sb_count)
                for sisi in site_count:
                    if wi==0: label =  'stie'+str(sisi)
                    else:label=None
                    locals()['y%s'%sisi] = Bintable_analysis_EDA[(Bintable_analysis_EDA['ELEMENT']==Diff_item)]['SITE'+str(sisi)+'_per'].iloc[1:x_axis_sb_count+1]
                    if y_bottom.sum()==0: plt.bar(x+(wi-1)*width,locals()['y%s'%sisi],label = label,color=coll[sisi],width=width*0.9,align='center')
                    else : plt.bar(x+(wi-1)*width,locals()['y%s'%sisi],bottom=y_bottom,label = label,color=coll[sisi],width=width*0.9,align='center')
                    y_bottom = y_bottom + np.array( locals()['y%s'%sisi])
                for coor in range(x_axis_sb_count): plt.annotate('', xy=(x[coor]+(wi-1)*width,y_bottom[coor]),xytext=(x[coor]+(wi-1)*width,y_bottom[coor]+0.0001),arrowprops=dict(facecolor=coll[wi], shrink=0.5))
                plt.scatter(0,0,marker='v',label = Diff_item,color=coll[wi])
                wi = wi + 1
            plt.gca().yaxis.set_major_formatter(FuncFormatter(to_percent))
            plt.yticks(fontsize=10)
            plt.xticks(x,np.array(Bintable_histogram["BIN_NU_NA"].iloc[1:x_axis_sb_count+1]),rotation=15,fontsize=7)
            plt.grid(True, linestyle = "-.", linewidth = 0.5)
            plt.legend(loc="upper right")
            plt.savefig(abcpath+'_Bintable_analysis_Diff.png',dpi=300)   
            plt.show()
            plt.close('all')

    #輸出by Tester資料
    if By_what==1 and len(BAA.NODE_NAM.unique())>1:
        while os.path.isfile(abcpath+'_Bintable_analysis_Tester.png')!=True:
            plt.figure(figsize=(14,7))
            total_width, n = 0.8, len(BAA.NODE_NAM.unique())
            width = total_width / n
            wi = 0#x軸位置
            for Tester_item in BAA.NODE_NAM.unique():
                y_bottom = np.array([0]*x_axis_sb_count)
                for sisi in site_count:
                    if wi==0: label =  'stie'+str(sisi)
                    else:label=None
                    locals()['y%s'%sisi] = Bintable_analysis_EDA[(Bintable_analysis_EDA['ELEMENT']==Tester_item)]['SITE'+str(sisi)+'_per'].iloc[1:x_axis_sb_count+1]
                    if y_bottom.sum()==0: plt.bar(x+(wi-1)*width,locals()['y%s'%sisi],label = label,color=coll[sisi],width=width*0.9,align='center')
                    else : plt.bar(x+(wi-1)*width,locals()['y%s'%sisi],bottom=y_bottom,label = label,color=coll[sisi],width=width*0.9,align='center')
                    y_bottom = y_bottom + np.array( locals()['y%s'%sisi])
                for coor in range(x_axis_sb_count): plt.annotate('', xy=(x[coor]+(wi-1)*width,y_bottom[coor]),xytext=(x[coor]+(wi-1)*width,y_bottom[coor]+0.0001),arrowprops=dict(facecolor=coll[wi], shrink=0.5))
                plt.scatter(0,0,marker='v',label = Tester_item,color=coll[wi])
                wi = wi + 1
            plt.gca().yaxis.set_major_formatter(FuncFormatter(to_percent))
            plt.yticks(fontsize=10)
            plt.xticks(x,np.array(Bintable_histogram["BIN_NU_NA"].iloc[1:x_axis_sb_count+1]),rotation=15,fontsize=7)
            plt.grid(True, linestyle = "-.", linewidth = 0.5)
            plt.legend(loc="upper right")
            plt.savefig(abcpath+'_Bintable_analysis_Tester.png',dpi=300)   
            plt.show()
            plt.close('all')
    w.Status.setText('Status: Bin Table 分析完成..')

   
def clickall(obj1,obj2):
    try:w.summary.cellChanged.disconnect()#關閉 "計算使用者勾選報表table的張數與顆數" 功能連接
    except:pass
    if obj1.isChecked():
        for aaa in range(obj2.rowCount()):
            item = QTableWidgetItem()
            item.setCheckState(Qt.Checked)
            obj2.setItem(aaa,0,item)
        if obj2==w.item_su:
            for aaa in range(w.item_search_su.rowCount()):
                item = QTableWidgetItem()
                item.setCheckState(Qt.Checked)
                w.item_search_su.setItem(aaa,0,item)            
    else:
        for aaa in range(obj2.rowCount()):
            item = QTableWidgetItem()
            item.setCheckState(Qt.Unchecked)
            obj2.setItem(aaa,0,item)        
        if obj2==w.item_su:
            for aaa in range(w.item_search_su.rowCount()):
                item = QTableWidgetItem()
                item.setCheckState(Qt.Unchecked)
                w.item_search_su.setItem(aaa,0,item)
    w.summary.cellChanged.connect(Summary_Count)#重新開啟 "計算使用者勾選報表table的張數與顆數" 功能連接
    Summary_Count()#觸發 "計算使用者勾選報表table的張數與顆數" 功能連接

def hidetime():#時間的勾勾設定
    if w.Need_time.isChecked():
        w.dateEdit_1.setEnabled(1)
        w.dateEdit_2.setEnabled(1)
    else:
        w.dateEdit_1.setDisabled(1)
        w.dateEdit_2.setDisabled(1)

def check_by_only_one(obj1):#專門針對那些打勾後>哪些就不能勾 or 可以勾
    if obj1 in [w.By_Tester,w.By_diff,w.By_LB,w.By_lot]:#給By_xxx用
        if obj1.isChecked()==False:
            w.Need_Combine.setChecked(False)
            w.Need_Combine.setDisabled(1)
        else:
            w.Need_Combine.setEnabled(1)
        if obj1!=w.By_Tester:w.By_Tester.setChecked(False)
        if obj1!=w.By_diff:w.By_diff.setChecked(False)
        if obj1!=w.By_LB:w.By_LB.setChecked(False)
        if obj1!=w.By_lot:w.By_lot.setChecked(False)

        
def get_sb():#得到目前勾選的summary有哪些SB
    if w.summary.rowCount()==0:
        w.Status.setText('Status: 無任何報表')
        alert('你腦瓦特了?~~無任何報表啊~~按個屁')
        return
    w.comboBox_sb.clear()#先清空下拉是清單
    findata_need=[]
    for aaa in range(w.summary.rowCount()):
        if w.summary.item(aaa,0).checkState()==2:
            findata_need.append(w.summary.item(aaa,column_summary.index('Filename')).text())
    if len(findata_need)==0:
        w.Status.setText('Status: 無任何報表被勾選..')
        alert('你個二百五看清楚好嗎>>沒有任何報表被勾選!!')
        return
    #選擇哪一場>>中華\銅鑼
    Factory = str(w.factory.currentText())
    if Factory=='中華': httpm = '******'
    else: httpm = '******'
    try_max_time=3
    try_time = 0
    while try_time < try_max_time:
        try:
            rs=requests.session()
            js_search2='******'
            header3={'******'}
            js3=json.loads(js_search2)
            js3['files']=findata_need
            res3=rs.post(httpm+'/api/datalog/filter/test/distinctOpts',headers=header3,json=js3)
            sb_num = res3.json()
            w.comboBox_sb.addItem('None')
            for ii in sb_num['SOFT_BIN']:
                w.comboBox_sb.addItem(str(ii)) 
            break
        except:
           try_time = try_time + 1
           time.sleep(try_time)


def Change_page(obj1):#換頁
    if obj1 ==w.CP_summary_1 or obj1 ==w.CP_item_1 or obj1 ==w.CP_group_1 or obj1 ==w.CP_CID_1:
        w.stackedWidget_su_it.setCurrentIndex(0)
    elif obj1 ==w.CP_summary_2 or obj1 ==w.CP_item_2 or obj1 ==w.CP_group_2 or obj1 ==w.CP_CID_2:
        w.stackedWidget_su_it.setCurrentIndex(2)
    elif obj1 ==w.CP_summary_3 or obj1 ==w.CP_item_3 or obj1 ==w.CP_group_3 or obj1 ==w.CP_CID_3:
        w.stackedWidget_su_it.setCurrentIndex(1)
    elif obj1 ==w.CP_summary_5 or obj1 ==w.CP_item_5 or obj1 ==w.CP_group_5 or obj1 ==w.CP_CID_5:
        w.stackedWidget_su_it.setCurrentIndex(3)
    elif obj1==w.pushButton_ch_SQL:
        if w.stackedWidget.currentIndex()==0:w.stackedWidget.setCurrentIndex(1)
        elif w.stackedWidget.currentIndex()==1:w.stackedWidget.setCurrentIndex(0)

    
def Summary_Count():#計算勾選報表後的張數/顆數/良率
    Dlog_u,Dlog_y,su_count,ea_count,go_count,yi_count=0,0,0,0.0,0.0,''
    for row in range(w.summary.rowCount()):
        if w.summary.item(row,0).checkState()==2:
            Dlog_u = int(w.summary.item(row,column_summary.index('Dlog_Units')).text())
            Dlog_y = float(w.summary.item(row,column_summary.index('Dlog_Yield')).text()[:len(w.summary.item(row,column_summary.index('Dlog_Yield')).text())-1])/100
            su_count = su_count + 1
            ea_count = ea_count + Dlog_u
            go_count = go_count + Dlog_u*Dlog_y
    try:#因為把最後一個勾勾取消>>為error
        if len(str((go_count / ea_count)*100))>=4:yi_count = str((go_count / ea_count)*100)[:4]+'%'
        else:yi_count = str((go_count / ea_count)*100)+'%'
        w.label_Summary_Count.setText('報表張數: '+str(su_count)+' 總顆數: '+str(int(ea_count))+'  總良率: '+yi_count)
    except:w.label_Summary_Count.setText('報表張數: '+str(su_count)+' 總顆數: '+str(ea_count)+'  總良率: '+yi_count)


def SQL_condition_search(combo_object,SC_object):#搜尋SQL_1
    try:import pymysql
    except:
        alert('抱歉~您非VIP會員無法使用~~~您要課金嗎?')
        return
    SC_object.clear()#先清空所有選項
    if combo_object.currentText()!='':#當前按的SC(下)其COMBO(上)必須有東西
        SSQL = 'SELECT distinct '+str(combo_object.currentText())+' FROM eisall WHERE '
        for ii in range(len(SQL_title_combo)):
            if (SQL_title_combo[ii].currentText()!='' and SQL_title_SC[ii].currentText()!='') and SQL_title_combo[ii] != combo_object:
                SSQL = SSQL + str(SQL_title_combo[ii].currentText()) + '="' + str(SQL_title_SC[ii].currentText()) + '" AND '
        if SSQL == 'SELECT distinct '+str(combo_object.currentText())+' FROM eisall WHERE ':#也就是前面沒有任何條件已經被選出來
            SSQL = 'SELECT distinct '+str(combo_object.currentText())+' FROM eisall AND '#後面加AND是讓後面統一去掉最後的AND
        SSQL = SSQL[:-5]#去掉最後的' AND '
        conn  =  pymysql.connect (host = '10.5.240.24',user = 'test',passwd = "123456",db = 'eis',charset = 'utf8')
        cur  =  conn.cursor ()
        print(SSQL)
        cur.execute(SSQL)
        data = cur.fetchall()
        SC_object.addItem('')
        for row in data:SC_object.addItem(str(row[0]))


def SQL_search():#搜尋SQL_2
    w.SQL_tabel.setSortingEnabled(False)#搜尋前>>先讓table不可排序>>因為可排序後搜尋會導致有些格子為空值
    while w.SQL_tabel.rowCount() > 0: w.SQL_tabel.removeRow(w.SQL_tabel.rowCount()-1) #先清空TABLE
    try:import pymysql
    except:
        alert('抱歉~您非VIP會員無法使用~~~您要課金嗎?')
        return SQL_title
    title_col = []
    SSQL = 'SELECT * FROM eisall WHERE '
    SSQL_TITLE = 'EXPLAIN eisall'
    for ii in range(len(SQL_title_combo)):
        if SQL_title_combo[ii].currentText()!='' and SQL_title_SC[ii].currentText()!='':   
            SSQL = SSQL + str(SQL_title_combo[ii].currentText()) + '="' + str(SQL_title_SC[ii].currentText()) + '" AND '
    if SSQL == 'SELECT * FROM eisall WHERE ':return#表示甚麼條件都還沒篩選
    SSQL = SSQL[:-5]#去掉最後的' AND '
    conn  =  pymysql.connect (host = '10.5.240.24',user = 'test',passwd = "123456",db = 'eis',charset = 'utf8')
    cur  =  conn.cursor ()
    print(SSQL)
    #先放入title_col
    cur.execute(SSQL_TITLE)
    data = cur.fetchall()
    title_col.append('cl')
    for row in data:title_col.append(list(row)[0])
    #建置SQL的Table
    w.SQL_tabel.setColumnCount(len(title_col))#設定table有多少row/col
    w.SQL_tabel.setHorizontalHeaderLabels(title_col)
    w.item_su.resizeColumnsToContents()
    w.item_su.resizeRowsToContents()
    #開始放入資料
    cur.execute(SSQL)
    data = cur.fetchall()
    for xyz in data:
        row = w.SQL_tabel.rowCount()
        w.SQL_tabel.setRowCount(row + 1)
        item = QTableWidgetItem()
        item.setFlags(item.flags())
        item.setCheckState(Qt.Unchecked)
        w.SQL_tabel.setItem(row,0,item)
        for coo in range(len(title_col)):
            try:w.SQL_tabel.setItem(row,coo+1,QTableWidgetItem(str(list(xyz)[coo])))
            except:w.SQL_tabel.setItem(row,coo+1,QTableWidgetItem(""))
    cur.close()
    conn.close()
    w.SQL_tabel.setSortingEnabled(True)#搜尋後>>再讓table可排序

def output_excel(obj1):#匯出excel
    abc=datetime.now()
    abcpath='D:/python/'+str(abc)[0:4]+'-'+str(abc)[5:7]+'-'+str(abc)[8:10]+'-'+str(abc)[11:13]+str(abc)[14:16]+str(abc)[17:19]+'_'
    if obj1==w.CP_group_4:
        if w.table_group_summary.rowCount()==0:
            alert('Are you kidding me?!~沒有資料~~你要輸出個毛線')
            return
        df_tem=pd.DataFrame(columns=column_group_summary)
        for row in range(w.table_group_summary.rowCount()):
            temp_list=[]
            for col in range(w.table_group_summary.columnCount()):
                temp_list.append(w.table_group_summary.item(row,col).text())
            df_tem.loc[row]=temp_list
        df_tem.to_excel (abcpath+'_group_summary.xlsx', index = None, header=True)
    elif obj1==w.CP_summary_4:
        if w.summary.rowCount()==0:
            alert('Are you kidding me?!~沒有資料~~你要輸出個毛線')
            return
        df_tem=pd.DataFrame(columns=column_summary[1:len(column_summary)])
        for row in range(w.summary.rowCount()):
            temp_list=[]
            for col in range(1,w.summary.columnCount()):
                temp_list.append(w.summary.item(row,col).text())
            df_tem.loc[row]=temp_list
        df_tem.to_excel (abcpath+'_summary.xlsx', index = None, header=True)
    elif obj1==w.CP_item_4:
        if w.item_su.rowCount()==0:
            alert('Are you kidding me?!~沒有資料~~你要輸出個毛線')
            return
        df_tem=pd.DataFrame(columns=column[1:len(column)])
        for row in range(w.item_su.rowCount()):
            temp_list=[]
            for col in range(1,w.item_su.columnCount()):
                temp_list.append(w.item_su.item(row,col).text())
            df_tem.loc[row]=temp_list
        df_tem.to_excel (abcpath+'_item.xlsx', index = None, header=True) 
    elif obj1==w.CP_CID_4:
        if w.table_CID.rowCount()==0:
            alert('Are you kidding me?!~沒有資料~~你要輸出個毛線')
            return
        df_tem=pd.DataFrame(columns=column_CID)
        for row in range(w.table_CID.rowCount()):
            temp_list=[]
            for col in range(w.table_CID.columnCount()):
                temp_list.append(w.table_CID.item(row,col).text())
            df_tem.loc[row]=temp_list
        df_tem.to_excel (abcpath+'_ChipID.xlsx', index = None, header=True)

    alert('已經輸出完成')



def ExamExam():#自動考試
    try:Examtest.exam(w.lineEdit_exam.text())
    except:alert('中斷....')


def CID_clear():#清空table_CID
    while w.table_CID.rowCount() > 0: w.table_CID.removeRow(w.table_CID.rowCount()-1)

def CID_Rule(obj1):
    global CID_Rule_group_check,CID_Rule_group_ID,CID_Rule_group_X,CID_Rule_group_Y,httpm
    if obj1 in [w.lineEdit_CID_ID_item,w.lineEdit_CID_X_item,w.lineEdit_CID_Y_item]:
        iiu=int(w.combox_CID_Rule.currentText())
        CID_Rule_group_ID[iiu]=w.lineEdit_CID_ID_item.text()
        CID_Rule_group_X[iiu]=w.lineEdit_CID_X_item.text()   
        CID_Rule_group_Y[iiu]=w.lineEdit_CID_Y_item.text()
        if CID_Rule_group_check[iiu]:w.checkBox_CID.setChecked(True)
        else:w.checkBox_CID.setChecked(False)
    elif obj1==w.combox_CID_Rule:
        w.lineEdit_CID_ID_item.setText(str(CID_Rule_group_ID[int(w.combox_CID_Rule.currentText())]))
        w.lineEdit_CID_X_item.setText(str(CID_Rule_group_X[int(w.combox_CID_Rule.currentText())]))
        w.lineEdit_CID_Y_item.setText(str(CID_Rule_group_Y[int(w.combox_CID_Rule.currentText())]))
        if CID_Rule_group_check[int(w.combox_CID_Rule.currentText())]:w.checkBox_CID.setChecked(True)
        else:w.checkBox_CID.setChecked(False)
    elif obj1==w.checkBox_CID:
        iiu=int(w.combox_CID_Rule.currentText())
        CID_Rule_group_check[iiu]=w.checkBox_CID.isChecked()
            
def CID_search():#CID_搜尋
    global rs,AuthToken,summary_need,item_all,findata,httpm
    #CID_Rule_group_ID,CID_Rule_group_X,CID_Rule_group_Y
    #=================前置卡關
    if item_all=='' or summary_need=='':
        w.Status.setText('Statu:無任何資料可查(報表與Item資料)')
        alert('無任何資料可查(報表與Item資料)')
        return
    if w.lineEdit_CID_item_need.text()=='':
        w.Status.setText('Statu:無輸入要查的item')
        alert('沒輸入要查的item~~耍呆?')
        return        
    if w.lineEdit_CID_ID_num.text()==''and w.lineEdit_CID_X_num.text()==''and w.lineEdit_CID_Y_num.text()=='':
        w.Status.setText('Status: 無輸入要查的ID_X_Y')
        alert('你沒填入任何ID X Y啊!!!,智障啊~~還是當我神? 知道你想搜尋甚麼?')
        return
        
    #===========搜尋與確認需要查的item==================
    CID_item_search_list,CID_list_temp=[],[]
    CID_list_temp=(str(w.lineEdit_CID_item_need.text()).split(','))
    for ppap in CID_list_temp:#如果有人key"240-247">>要變成[240,241,242...247]
        if ppap.find("-")!=-1:
            start,end= int(ppap[:ppap.find("-")]),int(ppap[ppap.find("-")+1:len(ppap)])
            while start <=end:
                CID_item_search_list.append(start)
                start = start + 1
        else:CID_item_search_list.append(ppap)
        
    #查詢需要查的item在item_all裡面第幾個
    CID_item_num_need=[]
    for num in range(len(item_all)):
        for itii in CID_item_search_list:
            if str(item_all[num]['TEST_NUM']) == str(itii):
                CID_item_num_need.append(num)
    print('CID_item_num_need : '+str(CID_item_num_need))
    
    if len(CID_item_num_need)==0:
        w.Status.setText('Status: 搜尋中斷')
        alert('請確認 需要搜尋的item是否key錯or是否存在')
        return   

    #===========搜尋與確認需要查的item Rule (ID_X_Y)===================
    #CID_Rule_group_ID,CID_Rule_group_X,CID_Rule_group_Y
    for rule_index in range(10):
        if CID_Rule_group_ID[rule_index]=='' or CID_Rule_group_check[rule_index]==False:continue#最前面卡關 >>沒設Rule當然就直接跳過 沒打勾也跳過
        
        CID_Rule_ID,CID_Rule_X,CID_Rule_Y="NA","NA","NA"
        CID_need_ID,CID_need_X,CID_need_Y=w.lineEdit_CID_ID_num.text(),w.lineEdit_CID_X_num.text(),w.lineEdit_CID_Y_num.text()
        for num in range(len(item_all)):
            if str(item_all[num]['TEST_NUM'])==CID_Rule_group_ID[rule_index]:CID_Rule_ID=num
            if str(item_all[num]['TEST_NUM'])==CID_Rule_group_X[rule_index]:CID_Rule_X=num
            if str(item_all[num]['TEST_NUM'])==CID_Rule_group_Y[rule_index]:CID_Rule_Y=num
        print('CID_Rule_ID : '+str(CID_Rule_ID)+'  CID_Rule_X : '+str(CID_Rule_X)+'  CID_Rule_Y : '+str(CID_Rule_Y)+'\n')
        
        #=========================卡關
        if CID_Rule_ID=="NA" or CID_Rule_X=="NA" or CID_Rule_Y=="NA":
            print('第'+str(rule_index)+'組:Rule有error,需確認')
            w.Status.setText('Status: 第'+str(rule_index)+'組:Rule有error,需確認')
            continue

        #===========接收有座標的DATALOG
        FILELIST=[]
        for filee in range(len(summary_need)):#給Json檔用
            locals()['filee_%s'%filee]={"FILE_NAME":"","PART_TYP":"","START_T":"","KLOT_NO":"FT","KY_PROGRAM":None,"TEST_TYP":"NA","TEST_MOD":"NA","TEST_MOD_NO":None,"KY_HBIN":None,"NODE_NAM":"","LOT_ID":""}
            locals()['filee_%s'%filee]["FILE_NAME"]=summary_need[filee]["FILE_NAME"]
            locals()['filee_%s'%filee]["PART_TYP"]=summary_need[filee]["PART_TYP"]
            locals()['filee_%s'%filee]["START_T"]=summary_need[filee]["START_T"]
            locals()['filee_%s'%filee]["KLOT_NO"]=summary_need[filee]["KLOT_NO"]
            locals()['filee_%s'%filee]["NODE_NAM"]=summary_need[filee]["NODE_NAM"]
            locals()['filee_%s'%filee]["LOT_ID"]=summary_need[filee]["LOT_ID"]
            FILELIST.append(locals()['filee_%s'%filee])
            del locals()['filee_%s'%filee]
        #選擇哪一場>>中華\銅鑼
        Factory = str(w.factory.currentText())
        if Factory=='中華': httpm = '******'
        else: httpm = '******'
        try_time=0
        try_max_time=5#超過一定次數~~就放棄(有時程式或昱冠會卡住)
        fuckyouasshole=0
        while try_time<try_max_time:#至少download資料要先成功>>才能有後續動作
            try:
                time.sleep(0.5)
                js_CID1='''******'''
                jd1=json.loads(js_CID1)
                jd1['FILELIST']=FILELIST
                jd1['RULES'][0]['TEST_RULE'][0]['ID']['TEST_TXT']=item_all[CID_Rule_ID]['TEST_TXT']
                jd1['RULES'][0]['TEST_RULE'][0]['ID']['TEST_NUM']=item_all[CID_Rule_ID]['TEST_NUM']
                jd1['RULES'][0]['TEST_RULE'][0]['X']['TEST_TXT']=item_all[CID_Rule_X]['TEST_TXT']
                jd1['RULES'][0]['TEST_RULE'][0]['X']['TEST_NUM']=item_all[CID_Rule_X]['TEST_NUM']
                jd1['RULES'][0]['TEST_RULE'][0]['Y']['TEST_TXT']=item_all[CID_Rule_Y]['TEST_TXT']
                jd1['RULES'][0]['TEST_RULE'][0]['Y']['TEST_NUM']=item_all[CID_Rule_Y]['TEST_NUM']
                headerd1={'******'}
                resd1=rs.post(httpm+'/api/chipid/getwaferdata',headers=headerd1,json=jd1)
                datalog_CID1=resd1.json()
                fuckyouasshole = datalog_CID1['RESULT'][0]['DATA'][0]['HARD_BIN']
                break
            except MemoryError:
                w.Status.setText('Status: 昱冠回傳檔案太大..請減少資料量or單一item畫畫看')
                alert('昱冠回傳檔案太大..請減少資料量or單一item畫畫看')
                try_time = try_max_time  
            except Exception:
                try_time = try_time + 1
        if try_time>=try_max_time:
            w.Status.setText('Status: 接收昱冠資料失敗~~繼續下一組Rule')
            print('aa,接收昱冠資料失敗~~繼續下一組Rule')
            #alert('接收昱冠資料失敗~~繼續下一組Rule')
            try:del js_CID1,jd1,resd1,datalog_CID1#獲取資料失敗>>就全部變數消除
            except:pass
            continue
    
        #===========接收有Val的DATALOG
        for it_need in CID_item_num_need:
            #選擇哪一場>>中華\銅鑼
            Factory = str(w.factory.currentText())
            if Factory=='中華': httpm = '******'
            else: httpm = '******'
            try_time=0
            try_max_time=5#超過一定次數~~就放棄(有時程式或昱冠會卡住)
            fuckyouasshole=0
            while try_time<try_max_time:#至少download資料要先成功>>才能有後續動作
                try:
                    time.sleep(0.5)
                    js_CID2='''******'''
                    jd2=json.loads(js_CID2)
                    jd2['files']=findata
                    jd2['tests'][0]['TEST_TXT']=item_all[it_need]['TEST_TXT']
                    jd2['tests'][0]['TEST_NUM']=item_all[it_need]['TEST_NUM']
                    headerd2={'******'}
                    resd2=rs.post(httpm+'/api/datalog/filter/test/data',headers=headerd2,json=jd2)
                    datalog_CID2=resd2.json()
                    fuckyouasshole = datalog_CID2[0]['DATA'][0]['HB']
                    break
                except MemoryError:
                    w.Status.setText('Status: 昱冠回傳檔案太大..請減少資料量or單一item畫畫看')
                    alert('昱冠回傳檔案太大..請減少資料量or單一item畫畫看')
                    try_time = try_max_time  
                except Exception:
                    try_time = try_time + 1
            if try_time>=try_max_time:
                w.Status.setText('Status: 接收昱冠資料失敗~~繼續下一組Rule')
                print('bb,接收昱冠資料失敗~~繼續下一組Rule')
                #alert('接收昱冠資料失敗~~繼續下一組Rule')
                try:del js_CID1,jd1,resd1,datalog_CID1,js_CID2,jd2,resd2,datalog_CID2#獲取資料失敗>>就全部變數消除
                except:pass
                continue

            #==========比對需要的資訊並填入table

            w.table_CID.setSortingEnabled(False)#搜尋前>>先讓table不可排序>>因為可排序後搜尋會導致有些格子為空值

            #先看有幾個條件要搜
            CID_ID_Yes,CID_X_Yes,CID_Y_Yes=True,True,True
            try:CID_need_ID = int(CID_need_ID)
            except:CID_ID_Yes=False
            try:CID_need_X = int(CID_need_X)
            except:CID_X_Yes=False
            try:CID_need_Y = int(CID_need_Y)
            except:CID_Y_Yes=False
            Rule_chose_num = sum([CID_ID_Yes,CID_X_Yes,CID_Y_Yes])#True個數
            #print('Rule_chose_num='+str(Rule_chose_num))

            #==========先搜尋有座標的DATALOG>>datalog_CID1
            df_IDXY = pd.DataFrame(columns=column_CID)
            df_IDXY_index=0
            for IDXYsummary in datalog_CID1["RESULT"]:
                try:#有可能有報表全部都o/s就不會有座標了
                    for IDXY in IDXYsummary["DATA"]:
                        RCN=0#Rule_chose_num>>True個數
                        if IDXY["ID"]==CID_need_ID:RCN=RCN+1
                        if IDXY["X_COORD"]==CID_need_X:RCN=RCN+1
                        if IDXY["Y_COORD"]==CID_need_Y:RCN=RCN+1
                        if RCN==Rule_chose_num:
                            for su in summary_need:
                                if su["FILE_NAME"]==IDXYsummary["FILE_NAME"]:
                                    print(su["FILE_NAME"])
                                    #column_CID=["PART_ID","Val","Item_num","ID","X_COORD","Y_COORD","HARD_BIN","ROM_COD","LTL","UTL","SOFT_BIN","SITE_NUM","LOT_ID","FILE_NAME","JOB_NAM","NODE_NAM","UNITS","START_T","ID_Item","X_Item","Y_Item"]
                                    IDXY_temp=[""]*len(column_CID)
                                    for appa in range(len(column_CID)):
                                        if column_CID[appa] =="ID_Item":#看是以哪個ID的ITEM來搜尋的
                                            IDXY_temp[appa] = str(item_all[CID_Rule_ID]['TEST_NUM'])
                                        elif column_CID[appa] =="X_Item":#看是以哪個X的ITEM來搜尋的
                                            IDXY_temp[appa] = str(item_all[CID_Rule_X]['TEST_NUM'])
                                        elif column_CID[appa] =="Y_Item":#看是以哪個Y的ITEM來搜尋的
                                            IDXY_temp[appa] = str(item_all[CID_Rule_Y]['TEST_NUM'])
                                        elif column_CID[appa] =="Item_num":#看是查哪個ITEM的值
                                            IDXY_temp[appa] = str(item_all[it_need]['TEST_NUM'])
                                        elif column_CID[appa] in ["LTL","UTL"]:#看是要查的ITEM其Limit值
                                            #因為FT與EQC的LIMIT有時不一樣,所以不能用item_all >>要用datalog_CID2去取得
                                            for summ in datalog_CID2:
                                                if summ["_id"]["FILE_NAME"]==su["FILE_NAME"]:
                                                    try:
                                                        if column_CID[appa]=="LTL":IDXY_temp[appa] = str(summ["LO_LIMIT"])
                                                        if column_CID[appa]=="UTL":IDXY_temp[appa] = str(summ["HI_LIMIT"])
                                                    except:IDXY_temp[appa] = "NA"#如果都沒查到就填NA
                                        else:
                                            try:
                                                IDXY_temp[appa]=IDXY[column_CID[appa]]#看資訊是不是在DATA
                                            except:
                                                try:
                                                    IDXY_temp[appa]=su[column_CID[appa]]#看資訊是不是在summary_need
                                                except:
                                                    IDXY_temp[appa]="NA"#都沒有就是NA
                                    #print(IDXY_temp)
                                    df_IDXY.loc[df_IDXY_index]=IDXY_temp
                                    df_IDXY_index = df_IDXY_index + 1
                except:pass

            #==========卡關
            if df_IDXY_index==0:
                w.Status.setText('Status: 無搜尋到符合的資料~~繼續下一組Rule')
                print('cc,無搜尋到符合的資料~~繼續下一組Rule')
                #alert('無搜尋到符合的資料~~繼續下一組Rule')
                try:del js_CID1,jd1,resd1,datalog_CID1,js_CID2,jd2,resd2,datalog_CID2#獲取資料失敗>>就全部變數消除
                except:pass
                break
            #==========將df_IDXY裡面的Val填入>>datalog_CID2
            for IDXYsummary2 in datalog_CID2:
                if IDXYsummary2["_id"]["FILE_NAME"] in df_IDXY["FILE_NAME"].unique().tolist():
                    for IDXY2 in IDXYsummary2["DATA"]:
                        for df_PRT  in range(len(df_IDXY["PART_ID"].tolist())):
                            if df_IDXY["PART_ID"].tolist()[df_PRT] == IDXY2["PRT"] and df_IDXY["FILE_NAME"].tolist()[df_PRT] == IDXYsummary2["_id"]["FILE_NAME"]:
                                df_IDXY.loc[df_PRT]["VAL"] = IDXY2["VAL"]

            #==========將df_IDXY的資訊放入table
            for it in range(df_IDXY.shape[0]):
                row = w.table_CID.rowCount()
                w.table_CID.setRowCount(row + 1)       
                for xyz in range(len(column_CID)):
                    try:
                        #column_CID=["PART_ID","VAL","Item_num","ID","X_COORD","Y_COORD","HARD_BIN","SOFT_BIN","SITE_NUM","LOT_ID","FILE_NAME","JOB_NAM","NODE_NAM","UNITS","ROM_COD","START_T","ID_Item","X_Item","Y_Item"]
                        if str(column_CID[xyz]) in ["PART_ID","ID","X_COORD","Y_COORD","HARD_BIN","SOFT_BIN","SITE_NUM","UNITS"]:
                            nameItem = QTableWidgetItem()
                            nameItem.setData(Qt.DisplayRole,int(df_IDXY.loc[it][str(column_CID[xyz])]))
                            w.table_CID.setItem(row,xyz,nameItem)
                        elif str(column_CID[xyz])=='VAL':
                            nameItem = QTableWidgetItem()
                            nameItem.setData(Qt.DisplayRole,float(df_IDXY.loc[it][str(column_CID[xyz])]))
                            w.table_CID.setItem(row,xyz,nameItem)                            
                        else:
                            w.table_CID.setItem(row,xyz,QTableWidgetItem(df_IDXY.loc[it][str(column_CID[xyz])]))
                    except:w.table_CID.setItem(row,xyz,QTableWidgetItem("NoData"))
            w.table_CID.setSortingEnabled(True)#搜尋後>>再讓table可排序
            
    alert('Chip_ID 搜尋完成')
    w.Status.setText('Status: Chip_ID 搜尋完成')

def CID_list_all():#CID_全列
    global rs,AuthToken,summary_need,item_all,findata,httpm
    #=================前置卡關
    if item_all=='' or summary_need=='':
        w.Status.setText('Statu:無任何資料可查(報表與Item資料)')
        alert('無任何資料可查(報表與Item資料)')
        return
    if w.lineEdit_CID_item_need.text()=='':
        w.Status.setText('Statu:至少要隨意輸入一個item')
        alert('至少要隨意輸入一個item')
        return        
    #===========搜尋與確認需要查的item==================
    CID_item_search_list,CID_list_temp=[],[]
    CID_list_temp=(str(w.lineEdit_CID_item_need.text()).split(','))
    try:CID_item_search_list=[CID_list_temp[0]]#如果有人key"240-247,191002150">>變成[240]>>也就是只取第一個
    except:#如果CID_list_temp為[] >>也就是空值時
        w.Status.setText('Status: 搜尋中斷')
        alert('請確認 需要搜尋的item是否key錯or是否存在')
        return  
    CID_item_num_need=[]#查詢需要查的item在item_all裡面第幾個
    for num in range(len(item_all)):
        if str(item_all[num]['TEST_NUM']) == str(CID_item_search_list[0]):
            CID_item_num_need.append(num)
    print('CID_item_num_need : '+str(CID_item_num_need))
    
    if len(CID_item_num_need)==0:
        w.Status.setText('Status: 搜尋中斷')
        alert('請確認 需要搜尋的item是否key錯or是否存在')
        return   
    #===========搜尋與確認需要查的item Rule ===================
    df_IDXY_all = []#合併全部的 df_IDXY>>放所有資料的DataFrame
    for rule_index in range(10):
        if CID_Rule_group_ID[rule_index]=='' or  CID_Rule_group_check[rule_index]==False:continue#最前面卡關 >>沒設Rule當然就直接跳過 沒打勾也跳過
        CID_Rule_ID,CID_Rule_X,CID_Rule_Y="NA","NA","NA"
        for num in range(len(item_all)):
            if str(item_all[num]['TEST_NUM'])==CID_Rule_group_ID[rule_index]:CID_Rule_ID=num
            if str(item_all[num]['TEST_NUM'])==CID_Rule_group_X[rule_index]:CID_Rule_X=num
            if str(item_all[num]['TEST_NUM'])==CID_Rule_group_Y[rule_index]:CID_Rule_Y=num
        print('CID_Rule_ID : '+str(CID_Rule_ID)+'  CID_Rule_X : '+str(CID_Rule_X)+'  CID_Rule_Y : '+str(CID_Rule_Y)+'\n')

        #=========================卡關
        if CID_Rule_ID=="NA" or CID_Rule_X=="NA" or CID_Rule_Y=="NA":
            print('第'+str(rule_index)+'組:Rule有error,需確認')
            w.Status.setText('Status: 第'+str(rule_index)+'組:Rule有error,需確認')
            continue

        #===========接收有座標的DATALOG
        FILELIST=[]
        for filee in range(len(summary_need)):#給Json檔用
            locals()['filee_%s'%filee]={"FILE_NAME":"","PART_TYP":"","START_T":"","KLOT_NO":"FT","KY_PROGRAM":None,"TEST_TYP":"NA","TEST_MOD":"NA","TEST_MOD_NO":None,"KY_HBIN":None,"NODE_NAM":"","LOT_ID":""}
            locals()['filee_%s'%filee]["FILE_NAME"]=summary_need[filee]["FILE_NAME"]
            locals()['filee_%s'%filee]["PART_TYP"]=summary_need[filee]["PART_TYP"]
            locals()['filee_%s'%filee]["START_T"]=summary_need[filee]["START_T"]
            locals()['filee_%s'%filee]["KLOT_NO"]=summary_need[filee]["KLOT_NO"]
            locals()['filee_%s'%filee]["NODE_NAM"]=summary_need[filee]["NODE_NAM"]
            locals()['filee_%s'%filee]["LOT_ID"]=summary_need[filee]["LOT_ID"]
            FILELIST.append(locals()['filee_%s'%filee])
            del locals()['filee_%s'%filee]
        #選擇哪一場>>中華\銅鑼
        Factory = str(w.factory.currentText())
        if Factory=='中華': httpm = '******'
        else: httpm = '******'
        try_time=0
        try_max_time=5#超過一定次數~~就放棄(有時程式或昱冠會卡住)
        fuckyouasshole=0
        while try_time<try_max_time:#至少download資料要先成功>>才能有後續動作
            try:
                time.sleep(0.5)
                js_CID1='''******'''
                jd1=json.loads(js_CID1)
                jd1['FILELIST']=FILELIST
                jd1['RULES'][0]['TEST_RULE'][0]['ID']['TEST_TXT']=item_all[CID_Rule_ID]['TEST_TXT']
                jd1['RULES'][0]['TEST_RULE'][0]['ID']['TEST_NUM']=item_all[CID_Rule_ID]['TEST_NUM']
                jd1['RULES'][0]['TEST_RULE'][0]['X']['TEST_TXT']=item_all[CID_Rule_X]['TEST_TXT']
                jd1['RULES'][0]['TEST_RULE'][0]['X']['TEST_NUM']=item_all[CID_Rule_X]['TEST_NUM']
                jd1['RULES'][0]['TEST_RULE'][0]['Y']['TEST_TXT']=item_all[CID_Rule_Y]['TEST_TXT']
                jd1['RULES'][0]['TEST_RULE'][0]['Y']['TEST_NUM']=item_all[CID_Rule_Y]['TEST_NUM']
                headerd1={'******'}
                resd1=rs.post(httpm+'/api/chipid/getwaferdata',headers=headerd1,json=jd1)
                datalog_CID1=resd1.json()
                fuckyouasshole = datalog_CID1['RESULT'][0]['DATA'][0]['HARD_BIN']
                break
            except MemoryError:
                w.Status.setText('Status: 昱冠回傳檔案太大..請減少資料量or單一item畫畫看')
                alert('昱冠回傳檔案太大..請減少資料量or單一item畫畫看')
                try_time = try_max_time  
            except Exception:
                try_time = try_time + 1
        if try_time>=try_max_time:
            w.Status.setText('Status: 接收昱冠資料失敗~~繼續下一組Rule')
            print('aa,接收昱冠資料失敗~~繼續下一組Rule')
            #alert('接收昱冠資料失敗~~繼續下一組Rule')
            try:del js_CID1,jd1,resd1,datalog_CID1#獲取資料失敗>>就全部變數消除
            except:pass
            continue
            
        #==========先搜尋有座標的DATALOG>>datalog_CID1
        for IDXYsummary in datalog_CID1["RESULT"]:                        
            df_IDXY = pd.DataFrame()
            try:#有可能有報表全部都o/s就不會有座標了
                if len(IDXYsummary["DATA"])!=0:
                    df_IDXY = json_normalize(IDXYsummary["DATA"])
                    for su in summary_need:
                        if su["FILE_NAME"]==IDXYsummary["FILE_NAME"]:
                            df_IDXY["ROM_COD"] = su["ROM_COD"]
                    df_IDXY["LOT_ID"] = IDXYsummary["LOT_ID"]
                    df_IDXY["FILE_NAME"] = IDXYsummary["FILE_NAME"]
                    df_IDXY_all = df_IDXY_all + np.array(df_IDXY).tolist()
                    column_CID_list = df_IDXY.columns.values.tolist()#得到Table的title
            except:pass
    
    #將list轉成DataFrame
    df_IDXY_all = pd.DataFrame(df_IDXY_all,columns=column_CID_list)
    #刪除重複資料
    df_IDXY_all = df_IDXY_all.drop_duplicates(subset=['ID','PART_ID','X_COORD','Y_COORD','FILE_NAME'])
    df_IDXY_all = df_IDXY_all.reset_index(drop=True)

    #輸出成excel用的檔名
    abc=datetime.now()
    abcpath='D:/python/'+str(abc)[0:4]+'-'+str(abc)[5:7]+'-'+str(abc)[8:10]+'-'+str(abc)[11:13]+str(abc)[14:16]+str(abc)[17:19]+'_'

    #輸出重複的excel
    df_IDXY_all_repeat = df_IDXY_all[df_IDXY_all.duplicated(subset=['ID','X_COORD','Y_COORD','FILE_NAME'],keep = False)].reset_index(drop=True)
    if df_IDXY_all_repeat.shape[0]!=0:
       df_IDXY_all_repeat.to_excel (abcpath+'_CID_list_all_repeat.xlsx', index = None, header=True) 
    #輸出全部的excel
    df_IDXY_all.to_excel (abcpath+'_CID_list_all.xlsx', index = None, header=True)

    alert('Chip_ID 搜尋完成[全列>>只輸出，不顯示]')
    w.Status.setText('Status: Chip_ID 搜尋完成[全列>>只輸出，不顯示]')

  
app = QApplication(sys.argv)
w = loadUi('untitled.ui')
w.setWindowTitle('MTFK')#視窗標題
#初始值
item_all,summary_all,summary_need,findata= '','','',[]
CID_Rule_group_check,CID_Rule_group_ID,CID_Rule_group_X,CID_Rule_group_Y=[True]*10,['']*10,['']*10,['']*10

#建置選擇時間
w.dateEdit_1.setCalendarPopup(True)#以日历的形式选择时间
w.dateEdit_2.setCalendarPopup(True)#以日历的形式选择时间
w.dateEdit_1.setDisabled(1)
w.dateEdit_2.setDisabled(1)
w.dateEdit_1.setDateTime(QDateTime.currentDateTime())#設置時間為now
w.dateEdit_2.setDateTime(QDateTime.currentDateTime())#設置時間為now

#建置table_summary
column_summary_width=[20,80,100,40,20,40,40,70,120,60,60,100,120]
column_summary = ['cl','Lot_Id','Start_Data','Test_Step','CMOD_COD','Dlog_Units','Dlog_Yield','Tester','Filename','Fab_Lot_id','Handler','Load_Id','JOB_REV','Program']
#cool = ['LOT_ID','START_T','ROM_COD','CMOD_COD','UNITS','YIELD','NODE_NAM','FILE_NAME','ENG_ID','HAND_ID','LOAD_ID','JOB_NAM']
#w.summary.setDragDropMode(QAbstractItemView.DragDrop)
#w.summary.setDragEnabled(True)
#w.summary.setAcceptDrops(True)
#w.summary.setSelectionBehavior(QAbstractItemView.SelectColumns)#以欄為單位
#w.summary.setEditTriggers(QAbstractItemView.NoEditTriggers)#不能编辑内容

w.summary.setColumnCount(len(column_summary))#設定table有多少row/col
w.summary.setHorizontalHeaderLabels(column_summary)
w.summary.resizeColumnsToContents()
w.summary.resizeRowsToContents()
w.summary.setFont(QFont('Times New Roman'))#對所有的單元格都使用這種字型
w.summary.setEditTriggers(QAbstractItemView.NoEditTriggers);#表格變為禁止編輯
w.summary.horizontalHeader().setStyleSheet('QHeaderView::section{background:green}')#设置表头的背景色为绿色
w.summary.setSortingEnabled(True)#设置表头可以自动排序
for cool in range(len(column_summary_width)): w.summary.setColumnWidth(cool, column_summary_width[cool])#调整列宽


#建置table_item
column_width=[20,45,180,45,45,50,40,40,40,40,40,40,40,40,40,40,40]
column = ['cl','TEST_NUM','TEST_TXT','SEQ','FAILS','EXECS','LTL','UTL',"SITE0","SITE1","SITE2","SITE3","SITE4","SITE5","SITE6","SITE7","SITE8"]#def有的會用到column>>不可刪
w.item_su.setColumnCount(len(column))#設定table有多少row/col
w.item_su.setHorizontalHeaderLabels(column)
w.item_su.resizeColumnsToContents()
w.item_su.resizeRowsToContents()
w.item_su.setFont(QFont('Times New Roman'))
w.item_su.setEditTriggers(QAbstractItemView.NoEditTriggers);
w.item_su.horizontalHeader().setStyleSheet('QHeaderView::section{background:green}')
w.item_su.setSortingEnabled(True)#设置表头可以自动排序
for cool in range(len(column_width)): w.item_su.setColumnWidth(cool, column_width[cool])#调整列宽

#建置table_item_search
w.item_search_su.setColumnCount(len(column))#設定table有多少row/col
w.item_search_su.setHorizontalHeaderLabels(column)
w.item_search_su.resizeColumnsToContents()
w.item_search_su.resizeRowsToContents()
w.item_search_su.setFont(QFont('Times New Roman'))
w.item_search_su.setEditTriggers(QAbstractItemView.NoEditTriggers);
w.item_search_su.horizontalHeader().setStyleSheet('QHeaderView::section{background:green}')
w.item_search_su.setSortingEnabled(True)#设置表头可以自动排序
for cool in range(len(column_width)): w.item_search_su.setColumnWidth(cool, column_width[cool])#调整列宽

#建置table_group_summary
#column_group_summary_width=[20,80,100,40,20,40,40,70,120,60,60,100,120]
column_group_summary = ['Group','Time','TEST_TXT','TEST_NUM','SEQ','EXECS','FAILS','MEAN','CPK','MAX','MIN','LTL','UTL']
w.table_group_summary.setColumnCount(len(column_group_summary))#設定table有多少row/col
w.table_group_summary.setHorizontalHeaderLabels(column_group_summary)
w.table_group_summary.resizeColumnsToContents()
w.table_group_summary.resizeRowsToContents()
w.table_group_summary.setFont(QFont('Times New Roman'))#對所有的單元格都使用這種字型
w.table_group_summary.setEditTriggers(QAbstractItemView.NoEditTriggers);#表格變為禁止編輯
w.table_group_summary.horizontalHeader().setStyleSheet('QHeaderView::section{background:green}')#设置表头的背景色为绿色
w.table_group_summary.setSortingEnabled(True)#设置表头可以自动排序
#for cool in range(len(column_group_summary_width)): w.table_group_summary.setColumnWidth(cool, column_group_summary_width[cool])#调整列宽

#建置table_CID
#column_CID只有第一個PART_ID的名字與位置是不能動的
column_CID=["PART_ID","VAL","Item_num","ID","X_COORD","Y_COORD","HARD_BIN","ROM_COD","LTL","UTL","SOFT_BIN","SITE_NUM","LOT_ID","FILE_NAME","JOB_NAM","NODE_NAM","UNITS","START_T","ID_Item","X_Item","Y_Item"]
column_CID_list=["PART_ID","Item_num","ID","X_COORD","Y_COORD","HARD_BIN","LOT_ID","FILE_NAME","START_T","ID_Item","X_Item","Y_Item"]
w.table_CID.setColumnCount(len(column_CID))#設定table有多少row/col
w.table_CID.setHorizontalHeaderLabels(column_CID)#設定table第一排的欄位名
w.table_CID.resizeColumnsToContents()
w.table_CID.resizeRowsToContents()
w.table_CID.setFont(QFont('Times New Roman'))#對所有的單元格都使用這種字型
w.table_CID.setEditTriggers(QAbstractItemView.NoEditTriggers);#表格變為禁止編輯
w.table_CID.horizontalHeader().setStyleSheet('QHeaderView::section{background:green}')#设置表头的背景色为绿色
w.table_CID.setSortingEnabled(True)#设置表头可以自动排序

#按鈕設定
w.Search_summary.clicked.connect(search_summary)
w.Search_item.clicked.connect(search_item)
w.Search_scatter.clicked.connect(search_scatter)
w.summary_click_all.stateChanged.connect(lambda:clickall(w.summary_click_all,w.summary))
w.summary_item_su_all.stateChanged.connect(lambda:clickall(w.summary_item_su_all,w.item_su))
w.Need_time.stateChanged.connect(lambda:hidetime())
w.By_Tester.stateChanged.connect(lambda:check_by_only_one(w.By_Tester))
w.By_diff.stateChanged.connect(lambda:check_by_only_one(w.By_diff))
w.By_LB.stateChanged.connect(lambda:check_by_only_one(w.By_LB))
w.By_lot.stateChanged.connect(lambda:check_by_only_one(w.By_lot))
w.Summary_search.clicked.connect(summary_search)
w.Get_sb.clicked.connect(get_sb)
w.Binning_Analysis.clicked.connect(Binning_Analysis)
w.pushButton_Group_By.clicked.connect(search_group_summary)
w.CP_group_4.clicked.connect(lambda:output_excel(w.CP_group_4))
w.CP_summary_4.clicked.connect(lambda:output_excel(w.CP_summary_4))
w.CP_item_4.clicked.connect(lambda:output_excel(w.CP_item_4))
w.CP_CID_4.clicked.connect(lambda:output_excel(w.CP_CID_4))
w.Auto_Exam.clicked.connect(ExamExam)
w.pushButton_CID_Search_2.clicked.connect(CID_search)
w.pushButton_CID_Search_4.clicked.connect(CID_list_all)
w.CID_Table_clear.clicked.connect(CID_clear)
w.checkBox_CID.setChecked(True)
w.checkBox_CID.stateChanged.connect(lambda:CID_Rule(w.checkBox_CID))
for item in [w.CB_site0,w.CB_site1,w.CB_site2,w.CB_site3,w.CB_site4,w.CB_site5,w.CB_site6,w.CB_site7,w.CB_site8]:#Site1-8的選項
    item.setChecked(True)

#按鈕說明
w.Item_num.setToolTip("輸入要列出的item數量後~~\n請在執行一次'Search_item'!!")
w.Summary_search.setToolTip("搜尋完報表後~~再篩選報表條件\n切記:\n     (1)篩選是一直疊加的\n     (2)如要重新篩選,請選擇'清除所有篩選'")
w.Get_sb.setToolTip("搜選(有勾選的)報表的SB\n此欄與'Search_item'相關:\n  切記:一旦選擇某個SB後在按'Search_item'則如同增加條件為某SB之搜尋條件")
w.Binning_Analysis.setToolTip("就是昱冠的Binning_Analysis")
w.Auto_Exam.setToolTip("只需填入課程代碼(ex:CO19020735),自動幫你閱讀與考試\nPS:只針對'能看錯誤題目'之考試才有效")
w.Search_item.setToolTip("")
w.Search_scatter.setToolTip("")
w.pushButton_Group_By.setToolTip("")
w.checkBox_PPT.setToolTip("")
w.Need_Histogram.setToolTip("")
w.ONSEMI_note.setToolTip("")
w.By_Tester.setToolTip("")
w.By_diff.setToolTip("")
w.By_LB.setToolTip("")
w.By_lot.setToolTip("")
w.Need_Combine.setToolTip("")
w.pushButton_CID_Search_2.setToolTip("MTFK")

#第一次開啟程式時>>有些按鈕要設為不可按
w.Search_scatter.setDisabled(1)
w.Need_Combine.setDisabled(1)
w.Search_item.setDisabled(1)
w.pushButton_Group_By.setDisabled(1)
w.pushButton_CID_Search_2.setDisabled(1)
w.Status.setText('Status: 待操作')

#換頁按鈕
w.CP_group_1.clicked.connect(lambda:Change_page(w.CP_group_1))
w.CP_group_2.clicked.connect(lambda:Change_page(w.CP_group_2))
w.CP_group_3.clicked.connect(lambda:Change_page(w.CP_group_3))
w.CP_group_5.clicked.connect(lambda:Change_page(w.CP_group_5))
w.CP_item_1.clicked.connect(lambda:Change_page(w.CP_item_1))
w.CP_item_2.clicked.connect(lambda:Change_page(w.CP_item_2))
w.CP_item_3.clicked.connect(lambda:Change_page(w.CP_item_3))
w.CP_item_5.clicked.connect(lambda:Change_page(w.CP_item_5))
w.CP_summary_1.clicked.connect(lambda:Change_page(w.CP_summary_1))
w.CP_summary_2.clicked.connect(lambda:Change_page(w.CP_summary_2))
w.CP_summary_3.clicked.connect(lambda:Change_page(w.CP_summary_3))
w.CP_summary_5.clicked.connect(lambda:Change_page(w.CP_summary_5))
w.CP_CID_1.clicked.connect(lambda:Change_page(w.CP_CID_1))
w.CP_CID_2.clicked.connect(lambda:Change_page(w.CP_CID_2))
w.CP_CID_3.clicked.connect(lambda:Change_page(w.CP_CID_3))
w.CP_CID_5.clicked.connect(lambda:Change_page(w.CP_CID_5))
w.pushButton_ch_SQL.clicked.connect(lambda:Change_page(w.pushButton_ch_SQL))


#很讚!!>>小遊戲
w.you_cant_see_me.clicked.connect(xxx)

#輸入欄設置
all_lineEdit = [w.lineEdit_item,w.lineEdit_lot,w.lineEdit_file,w.lineEdit_pg,w.lineEdit_tester,w.lineEdit_cmod,w.lineEdit_summary,
                w.lineEdit_XMIN,w.lineEdit_XMAX,
                w.lineEdit_CID_ID_item,w.lineEdit_CID_X_item,w.lineEdit_CID_Y_item,
                w.lineEdit_CID_ID_num,w.lineEdit_CID_X_num,w.lineEdit_CID_Y_num,
                w.lineEdit_CID_item_need]
for lineEdit in all_lineEdit:
    lineEdit.setClearButtonEnabled(True)  
w.lineEdit_item.editingFinished.connect(item_search)
w.lineEdit_CID_ID_item.editingFinished.connect(lambda:CID_Rule(w.lineEdit_CID_ID_item))
w.lineEdit_CID_X_item.editingFinished.connect(lambda:CID_Rule(w.lineEdit_CID_X_item))
w.lineEdit_CID_Y_item.editingFinished.connect(lambda:CID_Rule(w.lineEdit_CID_Y_item))

#將畫面預設為轉到summary表
w.stackedWidget_su_it.setCurrentIndex(0)

#下拉式清單設定
#w.summary_title.setEditable(True)
w.summary_title.addItem('清除所有篩選')
for ii in column_summary:
    if ii!='cl':w.summary_title.addItem(ii)
w.comboBox_count.addItem('>=')
w.comboBox_count.addItem('<=')
w.comboBox_Byxx.addItem('0,By_都不要')
w.comboBox_Byxx.addItem('1,By_Tester')
w.comboBox_Byxx.addItem('2,By_diff')
w.comboBox_Byxx.addItem('3,By_LB')
w.comboBox_GroupBy.addItem('0,By_Tester')
w.comboBox_GroupBy.addItem('1,By_diff')
w.comboBox_GroupBy.addItem('2,By_LB')
w.comboBox_GroupBy.addItem('3,By_Lot')
w.comboBox_GroupBy.addItem('4,By_TP_Rev')
w.comboBox_GroupBy.addItem('5,By_Handler')
for iii in range(10):
    w.combox_CID_Rule.addItem(str(iii))
w.combox_CID_Rule.currentIndexChanged.connect(lambda:CID_Rule(w.combox_CID_Rule))
w.factory.addItem('中華')
w.factory.addItem('銅鑼')
w.factory.setCurrentIndex(0)

#SQL的設定
SQL_title = ['Tester','CustomerLotNo','ProductFamily','OPNo','LotNo','DateCode','DeviceNo','PackageType','TraceCode','DiffusionLot','OriginalLotNo','Week','Month','Quarter','CORule','Weeksunday','LoadBoardNo','HandlerNo']
SQL_title_combo = [w.STitle1,w.STitle2,w.STitle3,w.STitle4,w.STitle5]
SQL_title_SC = [w.SC1,w.SC2,w.SC3,w.SC4,w.SC5]
for iii in SQL_title_combo:#增加Title選項
    for ppp in SQL_title:iii.addItem(ppp)
for iii in SQL_title_SC:iii.setEditable(True)#讓選項可以用key的

#設定選項旁邊的小按鈕
w.SC_B1.clicked.connect(lambda:SQL_condition_search(w.STitle1,w.SC1))
w.SC_B2.clicked.connect(lambda:SQL_condition_search(w.STitle2,w.SC2))
w.SC_B3.clicked.connect(lambda:SQL_condition_search(w.STitle3,w.SC3))
w.SC_B4.clicked.connect(lambda:SQL_condition_search(w.STitle4,w.SC4))
w.SC_B5.clicked.connect(lambda:SQL_condition_search(w.STitle5,w.SC5))
#設定SQL搜尋的主鍵
w.pushButton_SQL.clicked.connect(SQL_search)

#設定Group_summary的圖
w.Mplwidget.canvas = FigureCanvas(Figure(figsize=(12,12)))
vertical_layout = QVBoxLayout()
vertical_layout.addWidget(w.Mplwidget.canvas)
w.Mplwidget.canvas.axes = w.Mplwidget.canvas.figure.add_axes([0.1, 0.12, 0.88, 0.86])# left, bottom, width, height (range 0 to 1)
w.Mplwidget.setLayout(vertical_layout)


#FOR W488 的預設
CID_Rule_group_ID[0]='171000013'
CID_Rule_group_X[0]='171000011'
CID_Rule_group_Y[0]='171000012'
CID_Rule_group_ID[1]='171000113'
CID_Rule_group_X[1]='171000111'
CID_Rule_group_Y[1]='171000112'
CID_Rule_group_ID[2]='171000213'
CID_Rule_group_X[2]='171000211'
CID_Rule_group_Y[2]='171000212'
CID_Rule_group_ID[3]='171009013'
CID_Rule_group_X[3]='171009011'
CID_Rule_group_Y[3]='171009012'
#FOR WREE 的預設
CID_Rule_group_ID[4]='171000009'
CID_Rule_group_X[4]='171000007'
CID_Rule_group_Y[4]='171000008'
w.lineEdit_CID_ID_item.setText("171000013")
w.lineEdit_CID_X_item.setText("171000011")
w.lineEdit_CID_Y_item.setText("171000012")

w.show()
app.exec_()
